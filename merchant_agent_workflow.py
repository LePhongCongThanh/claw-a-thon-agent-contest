from __future__ import annotations

import asyncio
import json
import os
import re
import shutil
from datetime import datetime
from html import unescape
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from urllib.parse import parse_qs, unquote, urlparse

import pandas as pd

try:
    import numpy as np
except Exception:  # numpy luôn đi kèm pandas, nhưng phòng hờ
    np = None  # type: ignore

import requests
from dotenv import load_dotenv
from openai import AsyncOpenAI
from agents import Agent, Runner, handoff, function_tool, ModelSettings
from agents.models.openai_chatcompletions import OpenAIChatCompletionsModel


BASE_DIR = Path(__file__).resolve().parent
OUTPUT_DIR = BASE_DIR / "output"
UPLOAD_DIR = BASE_DIR / "uploads"

# RAG indexer — lazy, không crash nếu chưa cài chromadb/sentence-transformers
try:
    import rag_indexer as rag
    _RAG_AVAILABLE = True
except Exception:
    rag = None  # type: ignore
    _RAG_AVAILABLE = False

REQUIRED_COLUMNS = ["Date", "Merchant", "SOF_Type", "Acq_Type", "TPV"]
GROUP_COLUMNS = ["Merchant", "SOF_Type", "Acq_Type"]
EXCEL_TRANSACTION_SHEETS = ["Input_Template", "Input", "Transactions", "Raw_Data", "Data"]
FALSEY_ENV_VALUES = {"0", "false", "no", "off", "disabled"}

DATA_PREPARATION_INSTRUCTIONS = """
Data preparation workflow.
- Accept CSV, JSON, or Excel workbook uploads.
- For Excel, find the transaction sheet by schema, preferring Input_Template.
- Required transaction fields are Date, Merchant, SOF_Type, Acquisition_Type or Acq_Type, and TPV.
- Clean dates and TPV, standardize Acquisition_Type internally to Acq_Type, and drop invalid rows.
- Aggregate by Merchant, SOF_Type, and Acq_Type.
- Calculate MTD_TPV from the first day of the as-of month through the as-of date.
- Use MoM_Analysis.Prev_Month_TPV as the historical baseline when the workbook provides it; otherwise calculate previous month TPV from uploaded transactions.
- Calculate MoM_Growth_% = (MTD_TPV - Prev_Month_TPV) / Prev_Month_TPV.
- Save arranged transactions, segment metrics, and JSON summary into the output folder.
- Archive every uploaded source file into the uploads folder with a timestamped filename.
- Save an input summary markdown file and an output metric markdown report with timestamped names.
- Do not create business recommendations; only produce trustworthy prepared data and metrics.
""".strip()

ANALYTICS_INSTRUCTIONS = """
Bạn là Zalopay Merchant Analytics Assistant — trợ lý phân tích hiệu suất thanh toán cho các merchant của Zalopay.

## Vai trò
Sau khi merchant ký kết với Zalopay, mọi giao dịch của họ đi qua cổng thanh toán Zalopay.
Nhiệm vụ: phân tích TPV (Total Payment Volume), tìm nguyên nhân tăng/giảm, và đề xuất hành động tối ưu PnL.

## Quy tắc giao tiếp
- Không tự giới thiệu trừ khi được hỏi. Nếu hỏi: "Tôi là trợ lý phân tích merchant của Zalopay."
- KHÔNG nhắc đến tên file, đường dẫn, CSV, JSON, output folder với user. Chỉ trình bày insights.
- KHÔNG emit tool-call syntax trong câu trả lời.

## QUY TẮC HIỂN THỊ (rất quan trọng — chat chỉ render Markdown thuần, KHÔNG render LaTeX/math)
- TUYỆT ĐỐI KHÔNG dùng cú pháp toán/LaTeX: không `$...$`, không `\\uparrow`, `\\downarrow`, `\\times`, `\\%`.
- Thể hiện tăng/giảm bằng ký tự Unicode trực tiếp ▲ ▼ hoặc ↑ ↓, hoặc bằng chữ "tăng"/"giảm".
- Số liệu viết thẳng: "19.2M", "−4.9%", "+9.5%". Dùng dấu % và +/− bình thường.
- Bảng dùng Markdown table chuẩn. Bold dùng **text**.
- **TIẾNG VIỆT PHẢI CÓ DẤU ĐẦY ĐỦ VÀ NHẤT QUÁN** trong TOÀN BỘ câu trả lời — cả tiêu đề, heading, bảng,
  ghi chú và tiêu đề báo cáo PDF. KHÔNG viết tiếng Việt không dấu ("Bao Cao" SAI → "Báo Cáo" ĐÚNG).

## TRÍCH NGUỒN (cite) — bắt buộc khi dùng web search / crawl / social / nghiên cứu
- Mỗi khi dùng thông tin từ web search, crawl website hay social media → CITE nguồn ở CUỐI câu trả lời.
- Định dạng GỌN: một mục "**Nguồn:**", mỗi link 1 dòng kèm độ tin cậy. Ví dụ:
  - pharmacity.vn — banner Fundiin (HIGH)
  - cafef.vn/... — tin khuyến mãi (MEDIUM)
- Rút gọn URL (domain + path ngắn), không dán URL dài. KHÔNG bịa link — chỉ cite link tool thật trả về.

## Trả lời LINH HOẠT — không ép khuôn
- Trả lời ĐÚNG điều user hỏi. Nếu user hỏi 1 ý, trả lời gọn 1 ý — đừng nhồi đủ 6 mục.
- Tự chọn cấu trúc phù hợp với câu hỏi và dữ liệu thực tế. Không bịa số; chỉ dùng số từ tool trả về.
- Khi user muốn báo cáo tổng quan, có thể dùng các mục gợi ý: Tổng quan TPV & xu hướng, Segment tăng mạnh,
  Segment giảm + chẩn đoán, Segment mới, Hành động đề xuất, Web research findings. Đây là GỢI Ý, không bắt buộc.
- Có thể gợi ý xuất PDF khi phù hợp (vd sau một báo cáo dài), không cần hỏi máy móc sau mỗi câu.

## Độ chi tiết của Summary / Tóm tắt
- Khi user yêu cầu "summary", "tóm tắt", "tổng hợp" mà KHÔNG nói ngắn → làm ĐẦY ĐỦ, CHI TIẾT.
  Dài cũng được: bao quát mọi insight quan trọng (số liệu, xu hướng, nguyên nhân, đề xuất, nguồn).
- Khi user yêu cầu NGẮN ("ngắn gọn", "tóm tắt nhanh", "3 ý chính", "1 đoạn"...) → làm ĐÚNG theo ý user, súc tích.
- Dù dài hay ngắn: TUYỆT ĐỐI KHÔNG bịa đặt số liệu/sự kiện để chiều ý user hay để báo cáo trông "đầy đủ" hơn.
  Chỉ trình bày dữ liệu THẬT từ tool/log/hội thoại. Nếu thiếu dữ liệu → nói rõ "chưa có dữ liệu", không phỏng đoán.

## Kiến thức tham khảo để chẩn đoán (khung gợi ý — áp dụng linh hoạt, không máy móc)
- **Organic giảm**: so YoY (seasonal thì monitor); nếu bất thường → competitor lấy share / campaign kết thúc /
  phốt social (Threads, Facebook) → counter-campaign, relaunch, hoặc fix UX + retention.
- **Paid giảm**: breakdown voucher → budget cắt (bình thường, xét ROI) / budget giữ mà giảm (chất lượng kém,
  test creative-targeting) / budget tăng mà vẫn giảm (saturation, đổi voucher).
- **QR giảm**: thường KHÔNG tự chẩn đoán được → escalate BIZ/Area Manager kèm MTD vs Previous, YoY, vùng ảnh hưởng.
- **BNPL giảm**: merchant tự ra trả góp riêng (banner Home Credit/Kredivo) → crawl website merchant.
- **Tăng trưởng**: organic → amplify + replicate (ít budget); paid → check ROI trước rồi mới scale, theo dõi ad fatigue.

## Web Research — CHỦ ĐỘNG tìm, KHÔNG hỏi xin phép
- "Internet" = báo chí/tin tức + mạng xã hội (Threads, Facebook, TikTok) + các bài post/bài viết công khai.
- Khi user hỏi về một merchant (nguyên nhân, hoạt động, phốt, khuyến mãi...) hoặc nói "tra cứu thêm",
  "tìm hiểu thêm" → HÃY TỰ ĐỘNG nghiên cứu web NGAY (handoff cho Research Agent). TUYỆT ĐỐI KHÔNG hỏi
  "bạn có muốn tôi tìm kiếm trên web không?" — cứ tìm rồi báo kết quả.
- Thử NHIỀU góc truy vấn trước khi kết luận "không tìm thấy": tên merchant + "tin tức"/"khuyến mãi"/
  "trả góp"/"phốt"/"review"/"Threads"/"Facebook". Chỉ nói "chưa có thông tin" sau khi đã thử vài truy vấn.
- Confidence: Website merchant = HIGH; Báo chí/news = MEDIUM; Social public posts = LOW-MEDIUM. Luôn cite URL.
""".strip()


def _slug(value: str) -> str:
    value = re.sub(r"[^A-Za-z0-9._-]+", "_", value.strip())
    return value.strip("._-") or "merchant_transactions"


def _normalize_header(value: str) -> str:
    return re.sub(r"[^a-z0-9]+", "_", str(value).strip().lower()).strip("_")


def _canonical_column_map(columns: List[str]) -> Dict[str, str]:
    aliases = {
        "date": "Date",
        "transaction_date": "Date",
        "trans_date": "Date",
        "txn_date": "Date",
        "merchant": "Merchant",
        "merchant_name": "Merchant",
        "merchant_id": "Merchant",
        "sof_type": "SOF_Type",
        "sof": "SOF_Type",
        "source_of_fund": "SOF_Type",
        "source_of_funds": "SOF_Type",
        "payment_method": "SOF_Type",
        "acq_type": "Acq_Type",
        "acquisition_type": "Acq_Type",
        "acquisition": "Acq_Type",
        "channel": "Acq_Type",
        "tpv": "TPV",
        "amount": "TPV",
        "transaction_amount": "TPV",
        "payment_volume": "TPV",
        "total_payment_volume": "TPV",
    }
    mapping: Dict[str, str] = {}
    used: set = set()
    for column in columns:
        canonical = aliases.get(_normalize_header(column))
        if canonical and canonical not in used:
            mapping[column] = canonical
            used.add(canonical)
    return mapping


def _canonical_metric_column_map(columns: List[str]) -> Dict[str, str]:
    aliases = {
        "merchant": "Merchant",
        "merchant_name": "Merchant",
        "merchant_id": "Merchant",
        "sof_type": "SOF_Type",
        "sof": "SOF_Type",
        "source_of_fund": "SOF_Type",
        "source_of_funds": "SOF_Type",
        "payment_method": "SOF_Type",
        "acq_type": "Acq_Type",
        "acquisition_type": "Acq_Type",
        "acquisition": "Acq_Type",
        "channel": "Acq_Type",
        "mtd_tpv": "MTD_TPV",
        "month_to_date_tpv": "MTD_TPV",
        "prev_month_tpv": "Prev_Month_TPV",
        "previous_month_tpv": "Prev_Month_TPV",
        "mom_growth": "MoM_Growth_%",
        "mom_growth_pct": "MoM_Growth_%",
        "mom_growth_percent": "MoM_Growth_%",
    }
    mapping: Dict[str, str] = {}
    used: set = set()
    for column in columns:
        canonical = aliases.get(_normalize_header(column))
        if canonical and canonical not in used:
            mapping[column] = canonical
            used.add(canonical)
    return mapping


def _read_excel_transactions(path: Path) -> pd.DataFrame:
    excel = pd.ExcelFile(path)
    candidates = [sheet for sheet in EXCEL_TRANSACTION_SHEETS if sheet in excel.sheet_names]
    candidates.extend(sheet for sheet in excel.sheet_names if sheet not in candidates)

    for sheet in candidates:
        preview = pd.read_excel(excel, sheet_name=sheet, nrows=0)
        mapping = _canonical_column_map([str(column) for column in preview.columns])
        if all(column in mapping.values() for column in REQUIRED_COLUMNS):
            frame = pd.read_excel(excel, sheet_name=sheet)
            frame.attrs["source_sheet"] = sheet
            return frame

    raise ValueError(
        "No transaction sheet found. Expected a sheet like Input_Template with "
        "Date, Merchant, SOF_Type, Acquisition_Type, and TPV columns."
    )


def _read_reference_previous_month(path: Path) -> Tuple[Optional[pd.DataFrame], str]:
    if path.suffix.lower() not in {".xlsx", ".xlsm", ".xls"}:
        return None, "uploaded transactions"

    excel = pd.ExcelFile(path)
    if "MoM_Analysis" not in excel.sheet_names:
        return None, "uploaded transactions"

    reference = pd.read_excel(excel, sheet_name="MoM_Analysis")
    reference = reference.rename(
        columns=_canonical_metric_column_map([str(column) for column in reference.columns])
    )
    required = GROUP_COLUMNS + ["Prev_Month_TPV"]
    if any(column not in reference.columns for column in required):
        return None, "uploaded transactions"

    reference = reference[required].copy()
    reference["Prev_Month_TPV"] = pd.to_numeric(
        reference["Prev_Month_TPV"].astype(str).str.replace(r"[^\d.-]", "", regex=True),
        errors="coerce",
    )
    for column in GROUP_COLUMNS:
        reference[column] = reference[column].astype(str).str.strip()
    reference = reference.dropna(subset=required)
    reference = reference.groupby(GROUP_COLUMNS, dropna=False)["Prev_Month_TPV"].sum().reset_index()
    if reference.empty:
        return None, "uploaded transactions"
    return reference, "MoM_Analysis sheet"


def _read_uploaded_workbook_context(file_path: str) -> str:
    path = Path(file_path).expanduser().resolve()
    if path.suffix.lower() not in {".xlsx", ".xlsm", ".xls"}:
        return ""

    excel = pd.ExcelFile(path)
    context_blocks: List[str] = []
    for sheet_name in ["Detail_Analysis", "Voucher_Breakdown"]:
        if sheet_name not in excel.sheet_names:
            continue
        frame = pd.read_excel(excel, sheet_name=sheet_name).dropna(how="all")
        if frame.empty:
            continue
        records = frame.head(20).to_dict(orient="records")
        context_blocks.append(
            f"{sheet_name} data from uploaded workbook:\n"
            f"{json.dumps(records, ensure_ascii=False, default=str)}"
        )

    if not context_blocks:
        return ""
    return "\n\n".join(context_blocks)


def _archive_input_file(source_path: Path, timestamp: str) -> Path:
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    source_path = source_path.resolve()
    upload_root = UPLOAD_DIR.resolve()
    if upload_root in source_path.parents:
        return source_path

    suffix = source_path.suffix.lower()
    archived_path = UPLOAD_DIR / f"{timestamp}_{_slug(source_path.stem)}{suffix}"
    shutil.copy2(source_path, archived_path)
    return archived_path


def _workbook_sheet_names(path: Path) -> List[str]:
    if path.suffix.lower() not in {".xlsx", ".xlsm", ".xls"}:
        return []
    try:
        excel = pd.ExcelFile(path)
        return list(excel.sheet_names)
    except Exception:
        return []


def _markdown_table(records: List[Dict[str, Any]], columns: List[str]) -> str:
    if not records:
        return "_None._"

    header = "| " + " | ".join(columns) + " |"
    separator = "| " + " | ".join(["---"] * len(columns)) + " |"
    rows = []
    for record in records:
        values = []
        for column in columns:
            value = record.get(column, "")
            if isinstance(value, float):
                value = round(value, 2)
            values.append(str(value).replace("\n", " "))
        rows.append("| " + " | ".join(values) + " |")
    return "\n".join([header, separator, *rows])


def _write_input_summary_markdown(
    *,
    original_path: Path,
    archived_path: Path,
    frame: pd.DataFrame,
    source_sheet: str,
    timestamp: str,
) -> Path:
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    summary_path = UPLOAD_DIR / f"{timestamp}_{_slug(original_path.stem)}_input_summary.md"
    date_min = frame["Date"].min().strftime("%Y-%m-%d")
    date_max = frame["Date"].max().strftime("%Y-%m-%d")
    sheet_names = _workbook_sheet_names(archived_path)

    lines = [
        "# Input Archive Summary",
        "",
        f"- Created at: {timestamp}",
        f"- Original file: {original_path}",
        f"- Archived file: {archived_path}",
        f"- Source sheet: {source_sheet}",
        f"- Workbook sheets: {', '.join(sheet_names) if sheet_names else 'N/A'}",
        f"- Valid transaction rows: {len(frame)}",
        f"- Date range: {date_min} to {date_max}",
        f"- Merchant count: {frame['Merchant'].nunique()}",
        f"- SOF types: {', '.join(sorted(frame['SOF_Type'].dropna().unique()))}",
        f"- Acquisition types: {', '.join(sorted(frame['Acq_Type'].dropna().unique()))}",
        "",
        "## Purpose",
        "",
        "This archived file is retained so historical metrics can be recalculated later, "
        "including previous MTD or previous month baselines.",
    ]
    summary_path.write_text("\n".join(lines), encoding="utf-8")
    return summary_path


def _write_metrics_markdown_report(
    *,
    summary: Dict[str, Any],
    timestamp: str,
    stem: str,
) -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    report_path = OUTPUT_DIR / f"{timestamp}_{stem}_metrics_report.md"
    segment_columns = ["Merchant", "SOF_Type", "Acq_Type", "MTD_TPV", "Prev_Month_TPV", "MoM_Growth_%", "MoM_Status"]

    lines = [
        "# Merchant Metrics Report",
        "",
        f"- Created at: {timestamp}",
        f"- Source file: {summary['source_file']}",
        f"- Archived input file: {summary['archived_input_file']}",
        f"- Input summary file: {summary['input_summary_file']}",
        f"- Source sheet: {summary['source_sheet']}",
        f"- As-of date: {summary['as_of_date']}",
        f"- MTD period: {summary['mtd_period']}",
        f"- Previous month period: {summary['previous_month_period']}",
        f"- Previous month source: {summary['previous_month_source']}",
        f"- Rows ingested: {summary['rows_ingested']}",
        f"- Merchant count: {summary['merchant_count']}",
        f"- Segment count: {summary['segment_count']}",
        f"- Total MTD TPV: {round(summary['total_mtd_tpv'], 2)}",
        f"- Total previous month TPV: {round(summary['total_previous_month_tpv'], 2)}",
        "",
        "## Saved Data Files",
        "",
        f"- Arranged transactions: {summary['arranged_file']}",
        f"- Metrics CSV: {summary['metrics_file']}",
        f"- Metrics JSON: {summary['summary_file']}",
        "",
        "## Top High Growth Segments",
        "",
        _markdown_table(summary["top_high_growth_segments"], segment_columns),
        "",
        "## Top Underperforming Segments",
        "",
        _markdown_table(summary["top_underperforming_segments"], segment_columns),
        "",
        "## Top New Segments",
        "",
        _markdown_table(summary["top_new_segments"], segment_columns),
        "",
        "## Follow-up Usage Notes",
        "",
        "Use this report as the primary reference for follow-up questions. "
        "Use the archived input file only when recalculation is required.",
    ]
    report_path.write_text("\n".join(lines), encoding="utf-8")
    return report_path


def _load_saved_markdown_references(limit: int = 6, max_characters: int = 24000) -> str:
    files = []
    for folder in [OUTPUT_DIR, UPLOAD_DIR]:
        if folder.exists():
            files.extend(path for path in folder.glob("*.md") if path.is_file())
    files = sorted(files, key=lambda path: path.stat().st_mtime, reverse=True)[: max(1, limit)]

    blocks: List[str] = []
    remaining = max_characters
    for path in files:
        text = path.read_text(encoding="utf-8", errors="ignore")
        if remaining <= 0:
            break
        snippet = text[:remaining]
        remaining -= len(snippet)
        blocks.append(f"Reference markdown file: {path}\n{snippet}")

    if not blocks:
        return "No saved markdown references found yet."
    return "\n\n---\n\n".join(blocks)


def _env_flag(name: str, default: bool = True) -> bool:
    value = os.getenv(name)
    if value is None:
        return default
    return value.strip().lower() not in FALSEY_ENV_VALUES


def _strip_html(value: str) -> str:
    value = re.sub(r"<script.*?</script>", " ", value, flags=re.DOTALL | re.IGNORECASE)
    value = re.sub(r"<style.*?</style>", " ", value, flags=re.DOTALL | re.IGNORECASE)
    value = re.sub(r"<[^>]+>", " ", value)
    return re.sub(r"\s+", " ", unescape(value)).strip()


def _normalize_search_url(href: str) -> str:
    href = unescape(href or "").strip()
    if href.startswith("//"):
        href = f"https:{href}"
    elif href.startswith("/"):
        href = f"https://duckduckgo.com{href}"

    parsed = urlparse(href)
    query = parse_qs(parsed.query)
    redirected = query.get("uddg")
    if redirected:
        return unquote(redirected[0])
    return href


def _public_web_search(query: str, max_results: int = 4) -> List[Dict[str, str]]:
    query = query.strip()
    if not query:
        return []

    timeout = float(os.getenv("WEB_RESEARCH_TIMEOUT_SECONDS", "12"))
    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
            "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/125.0 Safari/537.36"
        )
    }
    endpoints = ["https://duckduckgo.com/html/", "https://html.duckduckgo.com/html/"]

    last_error: Optional[Exception] = None
    for endpoint in endpoints:
        try:
            response = requests.get(endpoint, params={"q": query}, headers=headers, timeout=timeout)
            response.raise_for_status()
            html = response.text
            break
        except requests.RequestException as exc:
            last_error = exc
    else:
        raise RuntimeError(f"Public web search failed for query '{query}': {last_error}")

    anchor_pattern = re.compile(
        r'<a[^>]+class="[^"]*result__a[^"]*"[^>]+href="(?P<href>[^"]+)"[^>]*>'
        r"(?P<title>.*?)</a>",
        flags=re.DOTALL | re.IGNORECASE,
    )
    anchors = list(anchor_pattern.finditer(html))
    results: List[Dict[str, str]] = []
    seen_urls: set[str] = set()

    for index, match in enumerate(anchors):
        title = _strip_html(match.group("title"))
        url = _normalize_search_url(match.group("href"))
        if not title or not url or url in seen_urls:
            continue
        seen_urls.add(url)

        next_start = anchors[index + 1].start() if index + 1 < len(anchors) else len(html)
        result_block = html[match.end() : next_start]
        snippet_match = re.search(
            r'<[^>]+class="[^"]*result__snippet[^"]*"[^>]*>(?P<snippet>.*?)</',
            result_block,
            flags=re.DOTALL | re.IGNORECASE,
        )
        snippet = _strip_html(snippet_match.group("snippet")) if snippet_match else ""
        results.append({"title": title, "url": url, "snippet": snippet})
        if len(results) >= max(1, int(max_results)):
            break

    return results


def _should_run_web_research(user_message: str, summaries: Optional[List[Dict[str, Any]]] = None) -> bool:
    if not _env_flag("WEB_RESEARCH_ENABLED", True):
        return False

    text = (user_message or "").lower()

    # Always run web research when user asks about a specific merchant by name
    merchant_inquiry_patterns = [
        re.compile(r"merchant\s+\w", re.IGNORECASE),
        re.compile(r"(thông tin|thong tin|information|info|nghiên cứu|nghien cuu|tìm hiểu|tim hieu)\s+(về|ve|about|on)\s+\w", re.IGNORECASE),
        re.compile(r"(về|ve|about)\s+(merchant|thương hiệu|thuong hieu|chuỗi|chuoi|brand)\s+\w", re.IGNORECASE),
        re.compile(r"cho\s+tôi\s+biết\s+.*(merchant|thương hiệu|chuỗi)", re.IGNORECASE),
        re.compile(r"(tell|show|give)\s+me\s+.*(merchant|about)", re.IGNORECASE),
    ]
    if any(p.search(user_message or "") for p in merchant_inquiry_patterns):
        return True

    trigger_terms = [
        "why", "reason", "root cause", "research", "search", "web", "social",
        "facebook", "threads", "competitor", "growth", "increase", "paid",
        "organic", "qr", "vietqr", "voucher", "campaign", "drop", "down",
        "decrease", "tăng", "tang", "tại sao", "tai sao", "nguyên nhân",
        "nguyen nhan", "lý do", "ly do", "nghiên cứu", "nghien cuu",
        "tìm hiểu", "tim hieu", "tìm kiếm", "tim kiem", "giảm", "giam",
        "tụt", "tut", "phốt", "phot", "đối thủ", "doi thu",
        "chiến dịch", "chien dich", "khuyến mãi", "khuyen mai",
        "thông tin", "thong tin", "cho biết", "cho biet",
    ]
    if any(term in text for term in trigger_terms):
        return True

    for summary in summaries or []:
        if (
            summary.get("top_underperforming_segments")
            or summary.get("top_high_growth_segments")
            or summary.get("top_new_segments")
        ):
            return True
    return False


def _research_section_priority(user_message: str) -> List[Tuple[str, str]]:
    text = (user_message or "").lower()
    growth_terms = ["growth", "up", "increase", "tăng", "tang", "cao", "high growth"]
    new_terms = ["new", "mới", "moi", "new segment"]
    drop_terms = ["drop", "down", "decrease", "giảm", "giam", "tụt", "tut", "underperform"]

    if any(term in text for term in new_terms):
        return [
            ("top_new_segments", "New segment diagnosis"),
            ("top_high_growth_segments", "Growth-up diagnosis"),
            ("top_underperforming_segments", "Decline diagnosis"),
        ]
    if any(term in text for term in growth_terms):
        return [
            ("top_high_growth_segments", "Growth-up diagnosis"),
            ("top_new_segments", "New segment diagnosis"),
            ("top_underperforming_segments", "Decline diagnosis"),
        ]
    if any(term in text for term in drop_terms):
        return [
            ("top_underperforming_segments", "Decline diagnosis"),
            ("top_high_growth_segments", "Growth-up diagnosis"),
            ("top_new_segments", "New segment diagnosis"),
        ]
    return [
        ("top_underperforming_segments", "Decline diagnosis"),
        ("top_high_growth_segments", "Growth-up diagnosis"),
        ("top_new_segments", "New segment diagnosis"),
    ]


def _research_targets_from_summaries(
    summaries: List[Dict[str, Any]],
    user_message: str = "",
    max_targets: int = 2,
) -> List[Dict[str, Any]]:
    targets: List[Dict[str, Any]] = []
    seen: set[Tuple[str, str, str]] = set()
    section_priority = _research_section_priority(user_message)
    for summary in summaries:
        for section_name, scenario in section_priority:
            for segment in summary.get(section_name, []):
                merchant = str(segment.get("Merchant", "")).strip()
                sof_type = str(segment.get("SOF_Type", "")).strip()
                acq_type = str(segment.get("Acq_Type", "")).strip()
                key = (merchant.lower(), sof_type.lower(), acq_type.lower())
                if not merchant or key in seen:
                    continue
                seen.add(key)
                targets.append(
                    {
                        "merchant": merchant,
                        "sof_type": sof_type,
                        "acq_type": acq_type,
                        "scenario": scenario,
                        "mom_growth": segment.get("MoM_Growth_%"),
                        "mtd_tpv": segment.get("MTD_TPV"),
                        "prev_month_tpv": segment.get("Prev_Month_TPV"),
                    }
                )
                if len(targets) >= max(1, int(max_targets)):
                    return targets
    return targets


def _extract_merchant_target_from_message(user_message: str) -> Optional[Dict[str, Any]]:
    text = re.sub(r"\s+", " ", user_message or "").strip()
    if not text:
        return None

    patterns = [
        r"(?:merchant|merchent)\s+(?P<merchant>[^\n,;]+)",
        r"(?:thông tin|thong tin|information|info)\s+(?:về|ve|about|on)\s+(?:merchant\s+)?(?P<merchant>[^\n,;]+)",
        r"(?:nghiên cứu|nghien cuu|tìm kiếm|tim kiem|tìm hiểu|tim hieu|research|search)\s+(?:merchant\s+)?(?P<merchant>[^\n,;]+)",
        r"(?:cho\s+(?:tôi|toi|me)\s+biết|tell\s+me|show\s+me)\s+.{0,30}?(?:merchant\s+|về\s+|ve\s+|about\s+)(?P<merchant>[^\n,;]+)",
        r"(?:về|ve|about)\s+(?P<merchant>[^\n,;]+)",
    ]
    stop_pattern = re.compile(
        r"\b(?:ở|o|kênh|kenh|channel|sof|drop|giảm|giam|tăng|tang|vì|vi|do|tại sao|tai sao)\b",
        flags=re.IGNORECASE,
    )
    for pattern in patterns:
        match = re.search(pattern, text, flags=re.IGNORECASE)
        if not match:
            continue
        merchant = match.group("merchant").strip(" .,:;!?\"'")
        stop_match = stop_pattern.search(merchant)
        if stop_match:
            merchant = merchant[: stop_match.start()].strip(" .,:;!?\"'")
        if len(merchant) >= 2:
            return {
                "merchant": merchant,
                "sof_type": "",
                "acq_type": "",
                "scenario": "Public merchant research",
                "mom_growth": None,
                "mtd_tpv": None,
                "prev_month_tpv": None,
            }
    return None


def _latest_metric_summaries(limit: int = 2) -> List[Dict[str, Any]]:
    if not OUTPUT_DIR.exists():
        return []

    summaries: List[Dict[str, Any]] = []
    paths = sorted(OUTPUT_DIR.glob("*_summary_*.json"), key=lambda item: item.stat().st_mtime, reverse=True)
    for path in paths[: max(1, int(limit))]:
        try:
            summaries.append(json.loads(path.read_text(encoding="utf-8")))
        except (OSError, json.JSONDecodeError):
            continue
    return summaries


def _append_unique(values: List[str], value: str, limit: int) -> None:
    cleaned = re.sub(r"\s+", " ", value).strip()
    if cleaned and cleaned not in values and len(values) < limit:
        values.append(cleaned)


def _build_public_research_queries(
    target: Optional[Dict[str, Any]],
    user_message: str,
    max_queries: int = 6,
) -> List[str]:
    queries: List[str] = []
    max_queries = max(1, int(max_queries))
    merchant = str((target or {}).get("merchant", "")).strip()
    sof_type = str((target or {}).get("sof_type", "")).strip()
    acq_type = str((target or {}).get("acq_type", "")).strip()
    scenario = str((target or {}).get("scenario", "")).strip().lower()
    message = re.sub(r"\s+", " ", user_message or "").strip()

    if merchant:
        quoted_merchant = f'"{merchant}"'
        acq_lower = acq_type.lower()
        sof_lower = sof_type.lower()

        if sof_type:
            _append_unique(queries, f"{quoted_merchant} {sof_type} thanh toán", max_queries)
        _append_unique(queries, f"{quoted_merchant} Zalopay thanh toán", max_queries)
        _append_unique(queries, f'site:facebook.com "{merchant}" thanh toán', max_queries)
        _append_unique(queries, f'site:threads.net "{merchant}"', max_queries)
        _append_unique(queries, f'site:tiktok.com "{merchant}"', max_queries)

        if acq_lower == "paid":
            _append_unique(queries, f"{quoted_merchant} voucher khuyến mãi thanh toán", max_queries)
            _append_unique(queries, f"{quoted_merchant} chiến dịch khuyến mãi", max_queries)
            _append_unique(queries, f"{quoted_merchant} quảng cáo ưu đãi thanh toán", max_queries)
        elif acq_lower == "organic":
            _append_unique(queries, f"{quoted_merchant} đối thủ thanh toán", max_queries)
            _append_unique(queries, f"{quoted_merchant} cộng đồng đánh giá", max_queries)

        if any(term in sof_lower for term in ["bnpl", "buy now", "pay later", "paylater", "trả góp", "tra gop"]):
            _append_unique(queries, f"{quoted_merchant} trả góp", max_queries)
            _append_unique(queries, f"{quoted_merchant} Buy Now Pay Later", max_queries)
            _append_unique(queries, f"{quoted_merchant} BNPL", max_queries)
        if any(term in sof_lower for term in ["qr", "vietqr"]):
            _append_unique(queries, f"{quoted_merchant} VietQR", max_queries)
            _append_unique(queries, f"{quoted_merchant} QR thanh toán lỗi", max_queries)
            _append_unique(queries, f"{quoted_merchant} cửa hàng thanh toán QR", max_queries)

        if "growth" in scenario or "new segment" in scenario:
            _append_unique(queries, f"{quoted_merchant} tăng trưởng khuyến mãi", max_queries)
            _append_unique(queries, f"{quoted_merchant} campaign thanh toán", max_queries)
            _append_unique(queries, f"{quoted_merchant} ra mắt dịch vụ mới", max_queries)

        _append_unique(queries, f"{quoted_merchant} thanh toán", max_queries)
        _append_unique(queries, f"{quoted_merchant} khuyến mãi voucher", max_queries)
        _append_unique(queries, f"{quoted_merchant} phốt khiếu nại", max_queries)
        _append_unique(queries, f"{quoted_merchant} lỗi thanh toán", max_queries)

    if not queries and message:
        compact_message = message[:160]
        _append_unique(queries, compact_message, max_queries)
        _append_unique(queries, f"{compact_message} Facebook", max_queries)
        _append_unique(queries, f"{compact_message} Threads", max_queries)
        _append_unique(queries, f"{compact_message} tin tức", max_queries)

    return queries


def _format_research_target(target: Optional[Dict[str, Any]]) -> str:
    if not target:
        return "User question"
    scenario = str(target.get("scenario", "")).strip()
    parts = [
        str(target.get("merchant", "")).strip(),
        str(target.get("sof_type", "")).strip(),
        str(target.get("acq_type", "")).strip(),
    ]
    label = " / ".join(part for part in parts if part)
    if scenario and label:
        label = f"{label} [{scenario}]"
    growth = target.get("mom_growth")
    if growth is not None:
        label = f"{label} (MoM growth: {growth}%)"
    return label or "User question"


def _write_public_research_markdown(
    *,
    timestamp: str,
    targets: List[Optional[Dict[str, Any]]],
    query_results: List[Dict[str, Any]],
) -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    report_path = OUTPUT_DIR / f"{timestamp}_public_web_research_report.md"

    lines = [
        "# Public Web Research Report",
        "",
        f"- Created at: {timestamp}",
        "- Scope: public web pages and social posts discoverable by a search engine.",
        "- Note: this does not access private, logged-in, or non-indexed Facebook/Threads content.",
        "",
        "## Research Targets",
        "",
    ]
    if targets:
        for target in targets:
            lines.append(f"- {_format_research_target(target)}")
    else:
        lines.append("- User question")

    for item in query_results:
        target_label = _format_research_target(item.get("target"))
        lines.extend(["", f"## Query: {item['query']}", "", f"- Target: {target_label}"])
        if item.get("error"):
            lines.append(f"- Search status: {item['error']}")
            continue
        results = item.get("results", [])
        if not results:
            lines.append("- No public search results found.")
            continue
        for result in results:
            lines.append(f"- [{result['title']}]({result['url']})")
            if result.get("snippet"):
                lines.append(f"  - Snippet: {result['snippet']}")

    report_path.write_text("\n".join(lines), encoding="utf-8")
    return report_path


def _resolve_merchant_website(merchant: str, timeout: float = 10.0) -> Optional[str]:
    """Search DuckDuckGo to find the official website URL of a merchant."""
    query = f'"{merchant}" trang web chính thức site'
    try:
        results = _public_web_search(query, max_results=3)
    except Exception:
        return None

    blocked_domains = {
        "facebook.com", "threads.net", "tiktok.com", "youtube.com",
        "twitter.com", "instagram.com", "zalo.me", "zalopay.vn",
        "duckduckgo.com", "google.com", "wikipedia.org",
    }
    merchant_slug = re.sub(r"[^a-z0-9]", "", merchant.lower())
    for result in results:
        url = result.get("url", "")
        try:
            parsed = urlparse(url)
            domain = parsed.netloc.lower().lstrip("www.")
            if any(blocked in domain for blocked in blocked_domains):
                continue
            domain_slug = re.sub(r"[^a-z0-9]", "", domain)
            if merchant_slug[:4] in domain_slug or domain_slug[:4] in merchant_slug:
                return f"{parsed.scheme}://{parsed.netloc}"
        except Exception:
            continue
    return None


def _fetch_page_text(url: str, timeout: float = 10.0, max_chars: int = 8000) -> str:
    """Fetch a webpage and return cleaned plain text."""
    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
            "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/125.0 Safari/537.36"
        ),
        "Accept-Language": "vi-VN,vi;q=0.9,en;q=0.8",
    }
    try:
        response = requests.get(url, headers=headers, timeout=timeout, allow_redirects=True)
        response.raise_for_status()
        return _strip_html(response.text)[:max_chars]
    except Exception:
        return ""


_BNPL_KEYWORDS = [
    "trả góp", "tra gop", "installment", "buy now pay later", "bnpl", "pay later",
    "paylater", "trả chậm", "tra cham", "góp 0%", "gop 0%", "lãi suất 0",
    "kredivo", "home credit", "mcredit", "fia credit", "aeon", "shinhan",
    "trả góp 0%", "mua trả góp",
]
_COMPETING_PAYMENT_KEYWORDS = [
    "momo", "mo mo", "shopee pay", "shopeepay", "vnpay", "vn pay",
    "grab pay", "grabpay", "moca", "airpay", "onepay", "payoo",
    "napas", "viettel pay", "viettelpay",
]
_SCANDAL_KEYWORDS = [
    "phốt", "phot", "lừa đảo", "lua dao", "khiếu nại", "khieu nai",
    "tố cáo", "to cao", "bóc phốt", "boc phot", "scam", "gian lận",
    "gian lan", "chặt chém", "chat chem", "thất vọng", "that vong",
]

_PAYMENT_SUBPATHS = [
    "/thanh-toan", "/phuong-thuc-thanh-toan", "/tra-gop",
    "/payment", "/checkout", "/installment",
]


def _detect_competing_services(text: str) -> Dict[str, List[str]]:
    """Scan page text for competing BNPL/payment keywords and return findings by category."""
    text_lower = text.lower()
    findings: Dict[str, List[str]] = {"bnpl_competitors": [], "payment_competitors": [], "scandal_signals": []}
    for kw in _BNPL_KEYWORDS:
        if kw in text_lower and kw not in findings["bnpl_competitors"]:
            findings["bnpl_competitors"].append(kw)
    for kw in _COMPETING_PAYMENT_KEYWORDS:
        if kw in text_lower and kw not in findings["payment_competitors"]:
            findings["payment_competitors"].append(kw)
    for kw in _SCANDAL_KEYWORDS:
        if kw in text_lower and kw not in findings["scandal_signals"]:
            findings["scandal_signals"].append(kw)
    return findings


def _crawl_merchant_website(merchant: str, sof_type: str = "") -> str:
    """Crawl the merchant's official website to detect competing payment/BNPL services.

    Returns a markdown-formatted summary of findings, or empty string if nothing found.
    """
    if not _env_flag("MERCHANT_WEBSITE_CRAWL_ENABLED", True):
        return ""

    timeout = float(os.getenv("WEB_RESEARCH_TIMEOUT_SECONDS", "12"))
    base_url = _resolve_merchant_website(merchant, timeout=timeout)
    if not base_url:
        return ""

    pages_to_check = [base_url]
    sof_lower = sof_type.lower()
    if any(t in sof_lower for t in ["bnpl", "buy now", "pay later", "paylater", "trả góp", "tra gop"]):
        for subpath in _PAYMENT_SUBPATHS:
            pages_to_check.append(base_url.rstrip("/") + subpath)
    else:
        pages_to_check.append(base_url.rstrip("/") + "/payment")
        pages_to_check.append(base_url.rstrip("/") + "/thanh-toan")

    all_findings: Dict[str, List[str]] = {"bnpl_competitors": [], "payment_competitors": [], "scandal_signals": []}
    crawled_urls: List[str] = []

    for url in pages_to_check[:4]:
        text = _fetch_page_text(url, timeout=timeout)
        if not text:
            continue
        crawled_urls.append(url)
        page_findings = _detect_competing_services(text)
        for category, keywords in page_findings.items():
            for kw in keywords:
                if kw not in all_findings[category]:
                    all_findings[category].append(kw)

    has_findings = any(all_findings[cat] for cat in all_findings)
    if not has_findings and not crawled_urls:
        return ""

    lines = [
        f"## Merchant Website Analysis: {merchant}",
        f"- Official website detected: {base_url}",
        f"- Pages crawled: {', '.join(crawled_urls) if crawled_urls else 'None reachable'}",
        "",
    ]
    if all_findings["bnpl_competitors"]:
        lines.append(
            "**BNPL / Installment services found on merchant site** "
            f"(potential reason for Zalopay BNPL drop): {', '.join(all_findings['bnpl_competitors'])}"
        )
    if all_findings["payment_competitors"]:
        lines.append(
            "**Competing payment methods found on merchant site**: "
            + ", ".join(all_findings["payment_competitors"])
        )
    if all_findings["scandal_signals"]:
        lines.append(
            "**Scandal/complaint signals detected on merchant site**: "
            + ", ".join(all_findings["scandal_signals"])
        )
    if not has_findings:
        lines.append("No competing payment services or scandal signals detected on the merchant's website.")

    return "\n".join(lines) + "\n"


def _run_public_web_research(user_message: str, summaries: Optional[List[Dict[str, Any]]] = None) -> str:
    summaries = summaries or []
    if not _should_run_web_research(user_message, summaries):
        return ""

    max_targets = int(os.getenv("WEB_RESEARCH_MAX_TARGETS", "2"))
    max_queries = int(os.getenv("WEB_RESEARCH_MAX_QUERIES_PER_TARGET", "10"))
    max_results = int(os.getenv("WEB_RESEARCH_MAX_RESULTS", "4"))
    # If user explicitly names a merchant, always prioritize researching that merchant first
    explicit_merchant_target = _extract_merchant_target_from_message(user_message)
    if explicit_merchant_target:
        targets: List[Optional[Dict[str, Any]]] = [explicit_merchant_target]
    else:
        targets = _research_targets_from_summaries(
            summaries,
            user_message=user_message,
            max_targets=max_targets,
        )
        if not targets:
            targets = [None]

    query_results: List[Dict[str, Any]] = []
    website_analysis_blocks: List[str] = []

    for target in targets:
        merchant = str((target or {}).get("merchant", "")).strip()
        sof_type = str((target or {}).get("sof_type", "")).strip()

        if merchant:
            website_block = _crawl_merchant_website(merchant, sof_type=sof_type)
            if website_block:
                website_analysis_blocks.append(website_block)

        queries = _build_public_research_queries(target, user_message, max_queries=max_queries)
        for query in queries:
            item: Dict[str, Any] = {"target": target, "query": query, "results": []}
            try:
                item["results"] = _public_web_search(query, max_results=max_results)
            except Exception as exc:
                item["error"] = str(exc)
            query_results.append(item)

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    report_path = _write_public_research_markdown(
        timestamp=timestamp,
        targets=targets,
        query_results=query_results,
    )
    report_text = report_path.read_text(encoding="utf-8", errors="ignore")

    output_parts = ["Public web research context. Use it as supporting evidence only, with citations and confidence levels:\n", report_text]
    if website_analysis_blocks:
        output_parts.append(
            "\n\nMerchant website crawl findings (direct evidence from merchant's own site — "
            "high confidence if keywords found):\n"
            + "\n\n".join(website_analysis_blocks)
        )
    return "\n".join(output_parts) + "\n\n"


def _read_transactions(file_path: str) -> pd.DataFrame:
    path = Path(file_path).expanduser().resolve()
    if not path.exists():
        raise FileNotFoundError(f"Uploaded file was not found: {path}")

    suffix = path.suffix.lower()
    if suffix == ".csv":
        frame = pd.read_csv(path)
    elif suffix in {".xlsx", ".xlsm"}:
        frame = _read_excel_transactions(path)
    elif suffix == ".xls":
        frame = _read_excel_transactions(path)
    elif suffix == ".json":
        frame = pd.read_json(path)
    else:
        raise ValueError("Upload a CSV, XLSX, XLSM, XLS, or JSON transaction file.")

    mapping = _canonical_column_map([str(column) for column in frame.columns])
    frame = frame.rename(columns=mapping)
    missing = [column for column in REQUIRED_COLUMNS if column not in frame.columns]
    if missing:
        raise ValueError(
            "Missing required columns: "
            + ", ".join(missing)
            + ". Expected Date, Merchant, SOF_Type, Acq_Type, TPV."
        )

    frame = frame[REQUIRED_COLUMNS].copy()
    frame["Date"] = pd.to_datetime(frame["Date"], errors="coerce")
    frame["TPV"] = pd.to_numeric(
        frame["TPV"].astype(str).str.replace(r"[^\d.-]", "", regex=True),
        errors="coerce",
    )
    for column in ["Merchant", "SOF_Type", "Acq_Type"]:
        frame[column] = frame[column].astype(str).str.strip()

    frame = frame.dropna(subset=["Date", "Merchant", "SOF_Type", "Acq_Type", "TPV"])
    frame = frame[frame["TPV"] >= 0]
    if frame.empty:
        raise ValueError("No valid transaction rows remained after cleaning the file.")

    if "source_sheet" not in frame.attrs:
        frame.attrs["source_sheet"] = "file"
    return frame.sort_values(["Date", "Merchant", "SOF_Type", "Acq_Type"]).reset_index(drop=True)


def _period_bounds(frame: pd.DataFrame, as_of_date: str) -> Tuple[pd.Timestamp, pd.Timestamp, pd.Timestamp, pd.Timestamp]:
    if as_of_date:
        as_of = pd.to_datetime(as_of_date, errors="raise")
    else:
        as_of = frame["Date"].max()

    month_start = pd.Timestamp(year=as_of.year, month=as_of.month, day=1)
    previous_month_end = month_start - pd.Timedelta(days=1)
    previous_month_start = pd.Timestamp(
        year=previous_month_end.year,
        month=previous_month_end.month,
        day=1,
    )
    return as_of.normalize(), month_start, previous_month_start, previous_month_end


def _growth_status(mtd_tpv: float, previous_tpv: float, growth_pct: Optional[float]) -> str:
    if previous_tpv == 0 and mtd_tpv > 0:
        return "New segment"
    if previous_tpv > 0 and mtd_tpv == 0:
        return "Dropped"
    if growth_pct is None:
        return "No comparable volume"
    if growth_pct >= 20:
        return "High growth"
    if growth_pct >= 0:
        return "Stable growth"
    return "Underperforming"


def calculate_merchant_metrics(file_path: str, as_of_date: str = "") -> str:
    """Arrange uploaded merchant transactions, calculate metrics, and save output files.

    Args:
        file_path: Absolute path of the uploaded transaction file.
        as_of_date: Optional analysis date in YYYY-MM-DD format. If blank, uses the latest
            transaction date in the uploaded file.
    """
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    original_source_path = Path(file_path).expanduser().resolve()
    source_path = _archive_input_file(original_source_path, timestamp)
    frame = _read_transactions(str(source_path))
    source_sheet = frame.attrs.get("source_sheet", "file")
    input_summary_path = _write_input_summary_markdown(
        original_path=original_source_path,
        archived_path=source_path,
        frame=frame,
        source_sheet=source_sheet,
        timestamp=timestamp,
    )
    as_of, month_start, prev_start, prev_end = _period_bounds(frame, as_of_date)

    arranged = frame.copy()
    arranged["Date"] = arranged["Date"].dt.strftime("%Y-%m-%d")
    arranged = arranged.rename(columns={"Acq_Type": "Acquisition_Type"})

    mtd_frame = frame[(frame["Date"] >= month_start) & (frame["Date"] <= as_of)]
    prev_frame = frame[(frame["Date"] >= prev_start) & (frame["Date"] <= prev_end)]

    mtd = (
        mtd_frame.groupby(GROUP_COLUMNS, dropna=False)["TPV"]
        .sum()
        .reset_index()
        .rename(columns={"TPV": "MTD_TPV"})
    )
    previous = (
        prev_frame.groupby(GROUP_COLUMNS, dropna=False)["TPV"]
        .sum()
        .reset_index()
        .rename(columns={"TPV": "Prev_Month_TPV"})
    )
    reference_previous, previous_month_source = _read_reference_previous_month(source_path)
    if reference_previous is not None:
        if previous.empty:
            previous = reference_previous
        else:
            previous = previous.merge(
                reference_previous,
                on=GROUP_COLUMNS,
                how="outer",
                suffixes=("_from_transactions", "_from_reference"),
            )
            previous["Prev_Month_TPV"] = previous["Prev_Month_TPV_from_reference"].combine_first(
                previous["Prev_Month_TPV_from_transactions"]
            )
            previous = previous[GROUP_COLUMNS + ["Prev_Month_TPV"]]
    metrics = mtd.merge(previous, on=GROUP_COLUMNS, how="outer").fillna(0)

    if metrics.empty:
        metrics = pd.DataFrame(columns=GROUP_COLUMNS + ["MTD_TPV", "Prev_Month_TPV"])

    metrics["MTD_TPV"] = metrics["MTD_TPV"].astype(float)
    metrics["Prev_Month_TPV"] = metrics["Prev_Month_TPV"].astype(float)
    metrics["MoM_Growth_%"] = metrics.apply(
        lambda row: None
        if row["Prev_Month_TPV"] == 0
        else round(((row["MTD_TPV"] - row["Prev_Month_TPV"]) / row["Prev_Month_TPV"]) * 100, 2),
        axis=1,
    )
    metrics["MoM_Status"] = metrics.apply(
        lambda row: _growth_status(row["MTD_TPV"], row["Prev_Month_TPV"], row["MoM_Growth_%"]),
        axis=1,
    )
    metrics["As_Of_Date"] = as_of.strftime("%Y-%m-%d")
    metrics["MTD_Period"] = f"{month_start:%Y-%m-%d} to {as_of:%Y-%m-%d}"
    metrics["Prev_Month_Period"] = f"{prev_start:%Y-%m-%d} to {prev_end:%Y-%m-%d}"
    metrics = metrics.sort_values(["MTD_TPV", "Prev_Month_TPV"], ascending=False).reset_index(drop=True)

    stem = _slug(original_source_path.stem)
    arranged_path = OUTPUT_DIR / f"{stem}_arranged_{timestamp}.csv"
    metrics_path = OUTPUT_DIR / f"{stem}_metrics_{timestamp}.csv"
    summary_path = OUTPUT_DIR / f"{stem}_summary_{timestamp}.json"

    arranged.to_csv(arranged_path, index=False)
    metrics.to_csv(metrics_path, index=False)

    comparable = metrics[metrics["Prev_Month_TPV"] > 0].copy()
    high_growth = comparable[comparable["MoM_Growth_%"] > 0].sort_values(
        "MoM_Growth_%",
        ascending=False,
    ).head(5)
    underperforming = comparable[comparable["MoM_Growth_%"] < 0].sort_values(
        "MoM_Growth_%",
        ascending=True,
    ).head(5)
    new_segments = metrics[metrics["MoM_Status"] == "New segment"].sort_values(
        "MTD_TPV",
        ascending=False,
    ).head(5)

    summary: Dict[str, Any] = {
        "source_file": str(original_source_path),
        "archived_input_file": str(source_path),
        "input_summary_file": str(input_summary_path),
        "source_sheet": source_sheet,
        "arranged_file": str(arranged_path),
        "metrics_file": str(metrics_path),
        "summary_file": str(summary_path),
        "previous_month_source": previous_month_source,
        "as_of_date": as_of.strftime("%Y-%m-%d"),
        "mtd_period": f"{month_start:%Y-%m-%d} to {as_of:%Y-%m-%d}",
        "previous_month_period": f"{prev_start:%Y-%m-%d} to {prev_end:%Y-%m-%d}",
        "rows_ingested": int(len(frame)),
        "merchant_count": int(frame["Merchant"].nunique()),
        "segment_count": int(len(metrics)),
        "total_mtd_tpv": float(metrics["MTD_TPV"].sum()),
        "total_previous_month_tpv": float(metrics["Prev_Month_TPV"].sum()),
        "top_high_growth_segments": high_growth.to_dict(orient="records"),
        "top_underperforming_segments": underperforming.to_dict(orient="records"),
        "top_new_segments": new_segments.to_dict(orient="records"),
    }
    markdown_report_path = _write_metrics_markdown_report(
        summary=summary,
        timestamp=timestamp,
        stem=stem,
    )
    summary["markdown_report_file"] = str(markdown_report_path)

    with summary_path.open("w", encoding="utf-8") as handle:
        json.dump(summary, handle, indent=2, ensure_ascii=False)

    return json.dumps(summary, indent=2, ensure_ascii=False)


# ---------------------------------------------------------------------------
# Flexible file inspection + dynamic computation (KHÔNG hardcode schema)
# Agent 1 dùng để tự xem cấu trúc file thật rồi viết pandas phù hợp.
# ---------------------------------------------------------------------------

def _load_any_file(file_path: str) -> Tuple[Dict[str, "pd.DataFrame"], "pd.DataFrame"]:
    """Load file bất kỳ định dạng. Trả về (sheets_dict, primary_df).

    - CSV/JSON: sheets = {"data": df}
    - Excel: sheets = {tên_sheet: df, ...} (đọc TẤT CẢ sheet, không ép sheet nào)
    """
    path = Path(file_path).expanduser().resolve()
    if not path.exists():
        raise FileNotFoundError(f"File không tồn tại: {path}")

    suffix = path.suffix.lower()
    if suffix == ".csv":
        df = pd.read_csv(path)
        return {"data": df}, df
    if suffix == ".tsv":
        df = pd.read_csv(path, sep="\t")
        return {"data": df}, df
    if suffix in {".xlsx", ".xlsm", ".xls"}:
        sheets = pd.read_excel(path, sheet_name=None)  # dict tất cả sheet
        primary = next(iter(sheets.values())) if sheets else pd.DataFrame()
        return sheets, primary
    if suffix == ".json":
        df = pd.read_json(path)
        return {"data": df}, df
    raise ValueError(f"Định dạng không hỗ trợ: {suffix}. Dùng CSV/TSV/XLSX/XLS/JSON.")


# Định dạng tài liệu (text/narrative) — đọc bằng extract_document_text, KHÔNG phải tabular
DOCUMENT_SUFFIXES = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".txt", ".md"}


def extract_document_text(file_path: str, max_chars: int = 20000) -> str:
    """Trích xuất text (và bảng) từ tài liệu PDF/Word/PowerPoint/text để agent đọc & phân tích.

    Returns text thuần. Bảng trong PDF được trích thành dạng 'a | b | c' theo dòng.
    """
    path = Path(file_path).expanduser().resolve()
    if not path.exists():
        raise FileNotFoundError(f"File không tồn tại: {path}")
    suffix = path.suffix.lower()
    parts: List[str] = []

    if suffix == ".pdf":
        try:
            import pdfplumber
            with pdfplumber.open(str(path)) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    txt = page.extract_text() or ""
                    if txt.strip():
                        parts.append(f"--- Trang {i} ---\n{txt}")
                    for tbl in page.extract_tables() or []:
                        rows = [" | ".join(str(c) if c is not None else "" for c in row) for row in tbl]
                        if rows:
                            parts.append("[Bảng]\n" + "\n".join(rows))
        except Exception:
            # Fallback: pypdf chỉ lấy text
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            for i, page in enumerate(reader.pages, 1):
                txt = page.extract_text() or ""
                if txt.strip():
                    parts.append(f"--- Trang {i} ---\n{txt}")

    elif suffix in {".docx", ".doc"}:
        import docx
        doc = docx.Document(str(path))
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for ti, table in enumerate(doc.tables, 1):
            rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
            if rows:
                parts.append(f"[Bảng {ti}]\n" + "\n".join(rows))

    elif suffix in {".pptx", ".ppt"}:
        from pptx import Presentation
        prs = Presentation(str(path))
        for si, slide in enumerate(prs.slides, 1):
            slide_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    slide_parts.append(shape.text_frame.text)
                if shape.has_table:
                    tbl = shape.table
                    rows = [" | ".join(c.text.strip() for c in row.cells) for row in tbl.rows]
                    slide_parts.append("[Bảng]\n" + "\n".join(rows))
            if slide_parts:
                parts.append(f"--- Slide {si} ---\n" + "\n".join(slide_parts))

    elif suffix in {".txt", ".md"}:
        parts.append(path.read_text(encoding="utf-8", errors="ignore"))

    else:
        raise ValueError(
            f"Định dạng tài liệu không hỗ trợ: {suffix}. "
            "Hỗ trợ: PDF, DOCX, PPTX, TXT, MD (dữ liệu bảng dùng CSV/XLSX)."
        )

    text = "\n\n".join(parts).strip()
    if not text:
        return "(Không trích xuất được nội dung text từ tài liệu — có thể là file scan/ảnh.)"
    return text[:max_chars]


def inspect_data_file(file_path: str, sample_rows: int = 5) -> str:
    """Trả về cấu trúc file (KHÔNG giả định schema): sheet, cột, dtype, mẫu dòng, thống kê.

    Dùng để Agent 1 hiểu file thật trước khi quyết định cách tính metrics.
    """
    sheets, _ = _load_any_file(file_path)
    report: Dict[str, Any] = {"file": str(Path(file_path).name), "sheets": {}}

    for name, df in sheets.items():
        cols_info = []
        for col in df.columns:
            series = df[col]
            info: Dict[str, Any] = {"name": str(col), "dtype": str(series.dtype)}
            non_null = series.dropna()
            # Gợi ý ngữ nghĩa: ngày, số, hay phân loại
            if pd.api.types.is_numeric_dtype(series):
                info["kind"] = "numeric"
                if len(non_null):
                    info["min"] = float(non_null.min())
                    info["max"] = float(non_null.max())
                    info["sum"] = float(non_null.sum())
            elif pd.api.types.is_datetime64_any_dtype(series):
                info["kind"] = "datetime"
                if len(non_null):
                    info["min"] = str(non_null.min())
                    info["max"] = str(non_null.max())
            else:
                info["kind"] = "categorical/text"
                uniques = non_null.astype(str).unique()
                info["n_unique"] = int(len(uniques))
                info["examples"] = [str(v) for v in uniques[:8]]
            info["n_null"] = int(series.isna().sum())
            cols_info.append(info)

        report["sheets"][name] = {
            "n_rows": int(len(df)),
            "n_cols": int(len(df.columns)),
            "columns": cols_info,
            "sample": df.head(sample_rows).astype(str).to_dict(orient="records"),
        }
    return json.dumps(report, ensure_ascii=False, indent=2, default=str)


# Builtins an toàn cho code execution (đủ xử lý pandas, chặn thao tác nguy hiểm như
# open/eval/exec/__import__/os). True/False/None là keyword nên không cần liệt kê.
import builtins as _builtins_module

_SAFE_BUILTIN_NAMES = (
    "abs", "all", "any", "bool", "dict", "divmod", "enumerate", "filter", "float",
    "format", "frozenset", "int", "isinstance", "issubclass", "len", "list", "map",
    "max", "min", "print", "range", "repr", "reversed", "round", "set", "slice",
    "sorted", "str", "sum", "tuple", "zip", "abs",
)
_SAFE_BUILTINS = {
    name: getattr(_builtins_module, name)
    for name in _SAFE_BUILTIN_NAMES
    if hasattr(_builtins_module, name)
}


def run_python_on_file(file_path: str, code: str, max_output_chars: int = 6000) -> str:
    """Chạy code pandas TUỲ Ý trên file đã upload và trả về output.

    Namespace có sẵn:
      - `pd`, `np`              : pandas, numpy
      - `df`                    : sheet/bảng chính (đã load)
      - `sheets`                : dict {tên_sheet: DataFrame} cho mọi sheet
      - `print(...)`            : in kết quả (được capture trả về)

    KHÔNG cần (và không được) import; pd/np đã sẵn. Không có truy cập os/file-system.
    """
    import io
    import contextlib

    sheets, primary = _load_any_file(file_path)
    namespace: Dict[str, Any] = {
        "pd": pd,
        "np": np,
        "df": primary,
        "sheets": sheets,
        "__builtins__": _SAFE_BUILTINS,
    }

    buffer = io.StringIO()
    try:
        with contextlib.redirect_stdout(buffer):
            exec(code, namespace)  # noqa: S102 — internal analytics tool, namespace bị giới hạn
    except Exception as exc:
        out = buffer.getvalue()
        return (out + f"\n[ERROR] {type(exc).__name__}: {exc}").strip()[:max_output_chars]

    out = buffer.getvalue().strip()
    # Nếu agent gán biến `result`, append nó vào output cho tiện
    if "result" in namespace and namespace["result"] is not None:
        result_val = namespace["result"]
        if isinstance(result_val, (pd.DataFrame, pd.Series)):
            out += "\n\n# result:\n" + result_val.to_string()
        else:
            out += f"\n\n# result: {result_val}"
    if not out:
        out = "(Code chạy xong nhưng không in gì. Dùng print(...) hoặc gán biến `result` để xem kết quả.)"
    return out[:max_output_chars]


def _safe_output_path(file_path: str) -> Path:
    candidate = Path(file_path)
    if not candidate.is_absolute():
        candidate = OUTPUT_DIR / candidate
    candidate = candidate.resolve()
    output_root = OUTPUT_DIR.resolve()
    if output_root not in candidate.parents and candidate != output_root:
        raise ValueError("Only files inside the output folder can be read.")
    if not candidate.exists():
        raise FileNotFoundError(f"Output file was not found: {candidate}")
    return candidate


def search_output_files(query: str = "", limit: int = 8) -> str:
    """Search generated output files by filename and small content snippets.

    Args:
        query: Optional text to search for in output filenames and file contents.
        limit: Maximum number of matching files to return.
    """
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    query_lower = query.lower().strip()
    matches: List[Dict[str, Any]] = []
    for path in sorted(OUTPUT_DIR.glob("*"), key=lambda item: item.stat().st_mtime, reverse=True):
        if not path.is_file():
            continue
        text = ""
        try:
            text = path.read_text(encoding="utf-8", errors="ignore")[:4000]
        except OSError:
            text = ""

        haystack = f"{path.name}\n{text}".lower()
        if query_lower and query_lower not in haystack:
            continue
        matches.append(
            {
                "file": str(path),
                "modified": datetime.fromtimestamp(path.stat().st_mtime).isoformat(timespec="seconds"),
                "snippet": text[:600],
            }
        )
        if len(matches) >= max(1, int(limit)):
            break

    if not matches:
        return "No matching output files found."
    return json.dumps(matches, indent=2, ensure_ascii=False)


def read_output_file(file_path: str, max_characters: int = 12000) -> str:
    """Read one generated output file from the output folder.

    Args:
        file_path: Output file path or filename.
        max_characters: Maximum number of characters to return.
    """
    path = _safe_output_path(file_path)
    text = path.read_text(encoding="utf-8", errors="ignore")
    return text[: max(1000, int(max_characters))]


def write_agent1_log(
    content: str,
    merchants: Optional[List[str]] = None,
    task_type: str = "analysis",
) -> str:
    """Write Agent 1 findings to a timestamped markdown log in the output folder.

    The log accumulates over time and serves as a historical knowledge base for
    Agent 2 to answer past-period questions (e.g. TPV trends, previous diagnoses).

    Args:
        content: Markdown content of the findings/analysis.
        merchants: List of merchant names involved (used in filename for easy search).
        task_type: Short label for the task: 'metrics', 'web_research', 'analysis', 'crawl'.
    """
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    merchant_slug = "_".join(_slug(m) for m in (merchants or []))[:40] or "general"
    filename = f"agent1_log_{timestamp}_{task_type}_{merchant_slug}.md"
    log_path = OUTPUT_DIR / filename

    header = (
        f"# Agent 1 Log — {task_type.title()}\n\n"
        f"- **Timestamp**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
        f"- **Task type**: {task_type}\n"
        f"- **Merchants**: {', '.join(merchants) if merchants else 'N/A'}\n\n"
        "---\n\n"
    )
    log_path.write_text(header + content, encoding="utf-8")

    # Auto-index vào ChromaDB ngay sau khi ghi
    if _RAG_AVAILABLE:
        try:
            rag.index_file(log_path)
        except Exception:
            pass

    return str(log_path)


def _md_inline_to_plain(text: str) -> str:
    """Gỡ markdown link [text](url) -> text. Giữ ** cho markdown=True của fpdf2 render bold."""
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = text.replace("`", "")
    return text


def _is_table_separator_row(cells: List[str]) -> bool:
    """Dòng phân cách markdown table kiểu |:---|:---:|---:|."""
    return bool(cells) and all(re.fullmatch(r":?-{1,}:?", c.strip()) or c.strip() == "" for c in cells)


def _parse_md_table(lines: List[str]) -> List[List[str]]:
    """Parse các dòng '| a | b |' thành list rows, bỏ dòng phân cách."""
    rows: List[List[str]] = []
    for line in lines:
        body = line.strip().strip("|")
        cells = [c.strip() for c in body.split("|")]
        if _is_table_separator_row(cells):
            continue
        rows.append(cells)
    # Chuẩn hóa số cột (pad cho đều)
    if rows:
        width = max(len(r) for r in rows)
        rows = [r + [""] * (width - len(r)) for r in rows]
    return rows


def export_analysis_pdf(analysis_text: str, title: str = "Merchant Analytics Report") -> str:
    """Generate a clean PDF report from markdown analysis text.

    Renders real tables, proper text wrapping, bold via markdown, and Vietnamese fonts.

    Args:
        analysis_text: Markdown analysis (headings, tables, bullets, bold).
        title: Report title shown at the top of the PDF.
    """
    from fpdf import FPDF
    from fpdf.enums import XPos, YPos
    from fpdf.fonts import FontFace

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    pdf_path = OUTPUT_DIR / f"merchant_analytics_report_{timestamp}.pdf"

    pdf = FPDF()
    pdf.set_margins(left=15, top=15, right=15)
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    # Unicode TTF font hỗ trợ tiếng Việt: bundled project font → system macOS.
    font_name = "Helvetica"
    _bundled = BASE_DIR / "static" / "fonts"
    _font_candidates = [
        (_bundled / "Arial.ttf", _bundled / "Arial-Bold.ttf"),
        (Path("/System/Library/Fonts/Supplemental/Arial.ttf"),
         Path("/System/Library/Fonts/Supplemental/Arial Bold.ttf")),
        (Path("/Library/Fonts/Arial Unicode.ttf"),
         Path("/Library/Fonts/Arial Unicode.ttf")),
    ]
    for regular, bold in _font_candidates:
        try:
            if regular.exists():
                pdf.add_font("Uni", "", str(regular))
                pdf.add_font("Uni", "B", str(bold if bold.exists() else regular))
                font_name = "Uni"
                break
        except Exception:
            continue

    unicode_font = font_name != "Helvetica"
    epw = pdf.epw  # effective page width (đã trừ lề)

    def safe(text: str) -> str:
        if unicode_font:
            return text
        return text.encode("latin-1", errors="replace").decode("latin-1")

    def write_block(text: str, size: int, style: str = "", gap: float = 1.5,
                    fill: bool = False, line_h: float = 5.5) -> None:
        pdf.set_font(font_name, style, size)
        pdf.multi_cell(
            epw, line_h, safe(_md_inline_to_plain(text)),
            markdown=True, fill=fill,
            new_x=XPos.LMARGIN, new_y=YPos.NEXT,  # FIX: luôn về lề trái dòng kế
        )
        if gap:
            pdf.ln(gap)

    def render_table(rows: List[List[str]]) -> None:
        if not rows:
            return
        pdf.set_font(font_name, "", 9)
        with pdf.table(
            text_align="LEFT",
            headings_style=FontFace(emphasis="BOLD", fill_color=(26, 79, 186), color=(255, 255, 255)),
            cell_fill_color=(245, 248, 255),
            cell_fill_mode="ROWS",
            line_height=6,
        ) as table:
            for row in rows:
                trow = table.row()
                for cell in row:
                    trow.cell(safe(_md_inline_to_plain(cell)))
        pdf.ln(2)

    # ── Title banner ────────────────────────────────────────────
    pdf.set_font(font_name, "B", 16)
    pdf.set_fill_color(26, 79, 186)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(0, 12, safe(title), new_x=XPos.LMARGIN, new_y=YPos.NEXT, fill=True, align="C")
    pdf.ln(2)
    pdf.set_font(font_name, "", 9)
    pdf.set_text_color(120, 120, 120)
    pdf.cell(0, 6, f"Generated: {datetime.now():%Y-%m-%d %H:%M:%S}  |  Zalopay Merchant Analytics",
             new_x=XPos.LMARGIN, new_y=YPos.NEXT, align="C")
    pdf.ln(4)
    pdf.set_text_color(0, 0, 0)

    # ── Body: parse markdown theo block ─────────────────────────
    lines = analysis_text.splitlines()
    i = 0
    while i < len(lines):
        stripped = lines[i].strip()

        if not stripped:
            pdf.ln(1.5)
            i += 1
            continue

        # Bỏ separator '---' và title cấp 1 (tránh trùng banner)
        if stripped == "---" or (stripped.startswith("# ") and not stripped.startswith("## ")):
            i += 1
            continue

        # Bảng: gom các dòng '|' liên tiếp
        if stripped.startswith("|"):
            block = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                block.append(lines[i])
                i += 1
            render_table(_parse_md_table(block))
            continue

        if stripped.startswith("## "):
            pdf.set_fill_color(230, 240, 255)
            write_block(stripped[3:], size=13, style="B", fill=True, gap=2, line_h=8)
        elif stripped.startswith("### "):
            write_block(stripped[4:], size=11, style="B", gap=1, line_h=6.5)
        elif stripped.startswith(("- ", "* ")):
            write_block("•  " + stripped[2:], size=10.5, gap=0.5)
        elif re.match(r"^\d+\.\s", stripped):
            write_block(stripped, size=10.5, gap=0.5)
        else:
            write_block(stripped, size=10.5, gap=1)
        i += 1

    pdf.output(str(pdf_path))
    return str(pdf_path)


# Phát hiện tiếng Việt KHÔNG dấu (heuristic): nhiều từ tiếng Việt hay gặp nhưng viết trần
_NO_DIACRITIC_HINTS = re.compile(
    r"\b(tang truong|sut giam|phan tich|danh gia|kien nghi|kenh|nguyen nhan|"
    r"toi uu|ra soat|kiem tra|day manh|thanh toan|tang|giam|hieu qua|"
    r"khuyen mai|van hanh|mo rong|nha thuoc|cong cu)\b",
    re.IGNORECASE,
)


def _likely_missing_diacritics(text: str) -> bool:
    """True nếu text có dấu hiệu tiếng Việt bị thiếu dấu (cần khôi phục)."""
    return bool(_NO_DIACRITIC_HINTS.search(text or ""))


async def restore_vietnamese_diacritics(text: str) -> str:
    """Khôi phục dấu tiếng Việt cho text (task hẹp → model làm tin cậy hơn là khi sinh báo cáo dài).
    GIỮ NGUYÊN markdown, số liệu, URL, cấu trúc — chỉ thêm dấu vào chữ tiếng Việt bị trần.
    """
    if not text or not _likely_missing_diacritics(text):
        return text  # đã đủ dấu → bỏ qua, không tốn LLM call
    api_key, base_url = _greennode_config()
    client = AsyncOpenAI(api_key=api_key, base_url=base_url)
    system_prompt = (
        "Bạn là công cụ KHÔI PHỤC DẤU tiếng Việt. Nhiệm vụ DUY NHẤT: thêm dấu tiếng Việt đầy đủ, "
        "chuẩn chính tả cho các từ bị viết trần (không dấu). "
        "TUYỆT ĐỐI GIỮ NGUYÊN: markdown (##, -, |, **), con số, %, ký hiệu, URL, xuống dòng, thứ tự, "
        "từ tiếng Anh và tên riêng. KHÔNG thêm/bớt/diễn giải nội dung. KHÔNG đổi số liệu. "
        "Chỉ trả về đúng văn bản đã thêm dấu, không thêm lời nào khác."
    )
    try:
        resp = await client.chat.completions.create(
            model=os.getenv("DIACRITIC_MODEL", AGENT2_MODEL),
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": text},
            ],
            temperature=0.0,
            max_tokens=int(os.getenv("DIACRITIC_MAX_TOKENS", "4000")),
        )
        fixed = (resp.choices[0].message.content or "").strip()
        # An toàn: chỉ nhận nếu kết quả không rỗng và độ dài hợp lý (tránh model cắt mất nội dung)
        if fixed and len(fixed) >= 0.6 * len(text):
            return fixed
    except Exception:
        pass
    return text  # lỗi → giữ nguyên


async def synthesize_conversation_report(conversation_text: str, focus: str = "") -> str:
    """Tổng hợp các Ý CHÍNH của cả cuộc hội thoại thành 1 báo cáo markdown sạch.

    Args:
        conversation_text: Toàn bộ transcript hội thoại (User/Assistant).
        focus: Nếu có, chỉ tổng hợp phần liên quan (vd 'Pharmacity', 'kênh BNPL').
    """
    api_key, base_url = _greennode_config()
    client = AsyncOpenAI(api_key=api_key, base_url=base_url)

    focus_line = (
        f"CHỈ tổng hợp nội dung liên quan đến: {focus}.\n" if focus.strip() else ""
    )
    system_prompt = (
        "Bạn là trợ lý tổng hợp báo cáo Zalopay Merchant Analytics. "
        "Nhiệm vụ: đọc cuộc hội thoại và TỔNG HỢP các ý chính thành MỘT báo cáo mạch lạc, "
        "không lặp, không chào hỏi, không hỏi lại. "
        "Định dạng Markdown sạch: dùng '## ' cho mục lớn, '- ' cho bullet, bảng markdown cho số liệu, "
        "'**...**' cho in đậm. TUYỆT ĐỐI KHÔNG dùng LaTeX/math ($...$, \\uparrow). "
        "TIẾNG VIỆT PHẢI CÓ DẤU ĐẦY ĐỦ VÀ NHẤT QUÁN — cả tiêu đề lẫn nội dung (vd 'Báo Cáo', KHÔNG 'Bao Cao'). "
        "Nếu hội thoại có trích dẫn nguồn (URL từ web/crawl/social), GIỮ LẠI trong mục '**Nguồn:**' ở cuối, "
        "rút gọn URL kèm độ tin cậy. "
        "Làm báo cáo ĐẦY ĐỦ, CHI TIẾT — dài cũng được, bao quát mọi insight quan trọng đã có trong hội thoại. "
        "TUYỆT ĐỐI KHÔNG bịa đặt số liệu/sự kiện để báo cáo trông đầy đủ hơn — chỉ dùng dữ liệu/insight ĐÃ "
        "xuất hiện trong hội thoại. Nếu thiếu dữ liệu thì ghi rõ, không phỏng đoán. "
        "Gợi ý cấu trúc: Tổng quan → Phân tích theo merchant/kênh → Nguyên nhân → Đề xuất hành động → Nguồn."
    )
    user_prompt = (
        f"{focus_line}"
        "Tổng hợp cuộc hội thoại sau thành một báo cáo merchant analytics hoàn chỉnh:\n\n"
        f"{conversation_text}"
    )
    try:
        resp = await client.chat.completions.create(
            model=AGENT2_MODEL,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            temperature=0.3,
            max_tokens=int(os.getenv("REPORT_SYNTH_MAX_TOKENS", "3500")),
        )
        content = (resp.choices[0].message.content or "").strip()
        content = content or conversation_text
        # Khôi phục dấu nếu model rớt dấu (đảm bảo PDF luôn có dấu nhất quán)
        content = await restore_vietnamese_diacritics(content)
        return content
    except Exception:
        # Fallback: nếu LLM lỗi, trả về transcript thô để vẫn xuất được PDF
        return conversation_text


# ---------------------------------------------------------------------------
# GreenNode config
# ---------------------------------------------------------------------------

def _greennode_config() -> Tuple[str, str]:
    """Return (api_key, base_url) for GreenNode — shared by both agents."""
    load_dotenv(BASE_DIR / ".env")
    api_key = (
        os.getenv("AI_PLATFORM_API_KEY")
        or os.getenv("OPENAI_API_KEY")
        or os.getenv("GREENNODE_API_KEY")
        or os.getenv("GREEENODE_API_KEY")
    )
    base_url = (
        os.getenv("OPENAI_BASE_URL")
        or os.getenv("GREENNODE_BASE_URL")
        or os.getenv("GREEENODE_BASE_URL")
        or "https://maas-llm-aiplatform-hcm.api.vngcloud.vn/v1"
    )
    if not api_key:
        raise RuntimeError("No API key found. Add GREEENODE_API_KEY to the .env file.")
    return api_key, base_url.rstrip("/")


def _make_client(model_name: str) -> OpenAIChatCompletionsModel:
    """Create an OpenAIChatCompletionsModel backed by the GreenNode endpoint."""
    api_key, base_url = _greennode_config()
    client = AsyncOpenAI(api_key=api_key, base_url=base_url)
    return OpenAIChatCompletionsModel(model=model_name, openai_client=client)


# ---------------------------------------------------------------------------
# Simple-message detection (fast path — skip agents for greetings)
# ---------------------------------------------------------------------------

_SIMPLE_MESSAGE_PATTERN = re.compile(
    r"^(hi|hello|hey|xin chào|chào|alo|ok|okay|cảm ơn|cam on|thanks|thank you|"
    r"bắt đầu|bat dau|start|help|giúp|giup|hướng dẫn|huong dan|\?+|\.+)$",
    re.IGNORECASE,
)
_ANALYTICAL_TERMS = re.compile(
    r"merchant|tpv|mom|mtd|growth|drop|giảm|tăng|phân tích|phan tich|"
    r"kênh|channel|organic|paid|qr|bnpl|voucher|segment|báo cáo|bao cao|"
    r"thông tin|thong tin|tìm|tim|nghiên|nghien|tại sao|tai sao|nguyên nhân",
    re.IGNORECASE,
)


def _is_simple_message(text: str) -> bool:
    text = text.strip()
    if _SIMPLE_MESSAGE_PATTERN.match(text):
        return True
    if len(text) <= 20 and not _ANALYTICAL_TERMS.search(text):
        return True
    return False


# ---------------------------------------------------------------------------
# Agent 1 — Research Agent  (google/gemma-4-31b-it)
# Tools: calculate metrics, web search, website crawl
# ---------------------------------------------------------------------------

AGENT1_MODEL = os.getenv("AGENT1_MODEL", "google/gemma-4-31b-it")

AGENT1_INSTRUCTIONS = """You are a Zalopay merchant analytics research agent (Agent 1).
Your job: gather data, compute metrics ACCURATELY from the real file, and ALWAYS write
findings to the persistent log.

## CRITICAL: never assume a fixed file schema

Uploaded files vary (different column names, layouts, wide/long formats, extra dimensions,
monthly columns, multiple sheets). NEVER guess columns. Always look at the real file first.

## Chọn tool theo LOẠI file
- **Tabular** (CSV, TSV, XLSX, XLS, JSON) → `inspect_file` rồi `run_python`/`calculate_merchant_metrics`.
- **Document** (PDF, Word .docx, PowerPoint .pptx, TXT, MD) → `read_document` để lấy text + bảng,
  rồi tự rút trích số liệu/insight từ nội dung đó. Nếu document chứa bảng số liệu, có thể tóm tắt/đối chiếu.
- Không chắc loại nào → thử `inspect_file` trước; nếu báo "định dạng không hỗ trợ" thì dùng `read_document`.

## Workflow khi có file tabular

1. **Inspect first** — call `inspect_file` to see actual sheets, columns, dtypes, value
   ranges and sample rows. Understand what the data really contains before computing.

2. **Choose the computation path:**
   - If the file clearly has Date + Merchant + SOF_Type + Acq_Type + TPV (or close aliases)
     → you MAY use `calculate_merchant_metrics` (fast path, also saves structured files).
   - OTHERWISE (any non-standard structure) → use `run_python` to write custom pandas:
     map the real columns to what you need, pick the right grouping, choose the period
     logic that fits the actual dates, and compute the metrics that make sense.
   - Decide thresholds/period windows from the DATA, not from fixed assumptions. Verify
     intermediate results with print() before concluding.

3. **Compute what the user actually asked for** — TPV totals, growth %, top/bottom segments,
   trends over the available periods, per-merchant or per-channel breakdowns, etc. Use
   `run_python` freely for any aggregation, pivot, or period-over-period calculation.

## Web research — CHỦ ĐỘNG, đa nguồn (KHÔNG hỏi xin phép)
   - `web_search` — "internet" gồm báo chí/tin tức + mạng xã hội (Threads, Facebook, TikTok) + bài post công khai.
     Khi cần thông tin ngoài về merchant → GỌI NGAY, thử NHIỀU truy vấn khác góc trước khi kết luận, ví dụ:
       "<merchant> tin tức", "<merchant> khuyến mãi", "<merchant> trả góp BNPL",
       "<merchant> phốt review", "<merchant> Threads", "<merchant> Facebook".
     Chỉ kết luận "chưa có thông tin" SAU KHI đã thử vài truy vấn mà đều rỗng.
   - `crawl_merchant_website` — phát hiện dịch vụ trả góp/BNPL/thanh toán đối thủ trên website merchant.
   - `search_history_files` / `read_history_file` — load past logs for trend comparison.

## Always at the end
4. **Write log** — ALWAYS call `write_log`:
   - Include: the actual numbers computed, how they were derived, diagnosis, recommendations
   - Set `task_type`: 'metrics', 'web_research', 'crawl', 'tpv_trend', or 'analysis'
   - Set `merchants` to comma-separated merchant names analyzed
5. **Return** a concise structured markdown report to Agent 2. Do not talk to the end user.

## Rules
- Never fabricate data. Report only what the tools actually return.
- If inspect shows the file lacks what's needed, say so clearly instead of guessing.
- When the user asks about historical TPV, search existing logs first before recomputing.
- **CITE NGUỒN**: khi dùng web_search / crawl_merchant_website / social, LUÔN kèm các URL nguồn
  (rút gọn) + độ tin cậy (HIGH/MEDIUM/LOW) trong report trả về Agent 2, để Agent 2 cite cho user.
- Viết tiếng Việt CÓ DẤU đầy đủ, nhất quán (không viết tiếng Việt không dấu).
"""


@function_tool
def agent1_inspect_file(file_path: str, sample_rows: int = 5) -> str:
    """Inspect ANY uploaded data file WITHOUT assuming a schema. Returns sheet names,
    column names + dtypes + semantic kind (numeric/datetime/categorical), value ranges,
    null counts, and sample rows. ALWAYS call this FIRST to understand the real structure
    before computing anything.

    Args:
        file_path: Absolute path to the uploaded file (CSV/TSV/XLSX/XLS/JSON).
        sample_rows: Number of sample rows to preview per sheet (default 5).
    """
    try:
        return inspect_data_file(file_path=file_path, sample_rows=sample_rows)
    except Exception as exc:
        return f"Inspect error: {type(exc).__name__}: {exc}"


@function_tool
def agent1_read_document(file_path: str, max_chars: int = 20000) -> str:
    """Read text + tables from a DOCUMENT file (PDF, Word .docx, PowerPoint .pptx, TXT, MD).
    Use this for non-tabular files — reports, slide decks, write-ups. For tabular data
    (CSV/XLSX) use inspect_file + run_python instead.

    Args:
        file_path: Absolute path to the document.
        max_chars: Max characters to return (default 20000).
    """
    try:
        return extract_document_text(file_path=file_path, max_chars=max_chars)
    except Exception as exc:
        return f"Read document error: {type(exc).__name__}: {exc}"


@function_tool
def agent1_run_python(file_path: str, code: str) -> str:
    """Run ARBITRARY pandas/numpy code on the uploaded file and return printed output.
    Use this for FLEXIBLE metric computation when the file does NOT match the standard
    schema, or when you need custom aggregations, period logic, pivots, or thresholds.

    Preloaded namespace (do NOT import anything):
      - pd, np         : pandas, numpy
      - df             : the primary sheet/table
      - sheets         : dict {sheet_name: DataFrame} for every sheet
      - print(...)     : capture output; or assign to `result` to auto-display

    Example:
      code = '''
      df["Date"] = pd.to_datetime(df["Date"])
      monthly = df.groupby(df["Date"].dt.to_period("M"))["TPV"].sum()
      print(monthly)
      result = monthly.pct_change()*100
      '''

    Args:
        file_path: Absolute path to the uploaded data file.
        code: Python code. pd/np/df/sheets are available; no imports, no file-system access.
    """
    try:
        return run_python_on_file(file_path=file_path, code=code)
    except Exception as exc:
        return f"Execution error: {type(exc).__name__}: {exc}"


@function_tool
def agent1_calculate_metrics(file_path: str, as_of_date: str = "") -> str:
    """STANDARD-SCHEMA fast path. Only use when the file already has columns
    Date, Merchant, SOF_Type, Acq_Type, TPV (or close aliases). Computes MTD TPV,
    Prev Month TPV, MoM growth and saves structured output/summary files used by the
    knowledge base. If inspect shows a different structure, use agent1_run_python instead.

    Args:
        file_path: Absolute path to the uploaded CSV or Excel file.
        as_of_date: Optional analysis date (YYYY-MM-DD). Defaults to the latest transaction date.
    """
    return calculate_merchant_metrics(file_path=file_path, as_of_date=as_of_date)


@function_tool
def agent1_web_search(query: str, max_results: int = 4) -> str:
    """Search the public internet (news sites, articles, and indexed social posts on
    Threads/Facebook/TikTok) for info about a merchant, campaign, voucher, incident, or market event.
    Call MULTIPLE times with different angles (news / promotion / installment / 'phốt review' /
    'Threads' / 'Facebook') before concluding nothing was found.

    Args:
        query: Search query in Vietnamese or English (vd: "Hasaki khuyến mãi", "Hasaki Threads phốt").
        max_results: Maximum number of results to return (default 4).
    """
    try:
        results = _public_web_search(query=query, max_results=max_results)
        return json.dumps(results, ensure_ascii=False, indent=2)
    except Exception as exc:
        return f"Search error: {exc}"


@function_tool
def agent1_crawl_merchant_website(merchant: str, sof_type: str = "") -> str:
    """Crawl the merchant's official website to detect competing BNPL/installment services or payment methods that may explain a Zalopay TPV drop.

    Args:
        merchant: Merchant name (e.g. 'Thế Giới Di Động', 'Long Châu').
        sof_type: Source of funds type (e.g. 'BNPL', 'VietQR') to focus the crawl.
    """
    return _crawl_merchant_website(merchant=merchant, sof_type=sof_type) or "No competing services detected."


@function_tool
def agent1_write_log(content: str, merchants: str = "", task_type: str = "analysis") -> str:
    """Save findings and analysis to a persistent markdown log file for future reference.
    ALWAYS call this at the end of every task to build the historical knowledge base.

    Args:
        content: Full markdown content of the findings — metrics, diagnosis, web research results, TPV trends.
        merchants: Comma-separated merchant names involved (e.g. 'Long Chau,The Gioi Di Dong').
        task_type: Task label — one of: 'metrics', 'web_research', 'crawl', 'tpv_trend', 'analysis'.
    """
    merchant_list = [m.strip() for m in merchants.split(",") if m.strip()]
    path = write_agent1_log(content=content, merchants=merchant_list, task_type=task_type)
    return f"Log saved: {path}"


def _build_agent1() -> Agent:
    return Agent(
        name="Research Agent",
        instructions=AGENT1_INSTRUCTIONS,
        model=_make_client(AGENT1_MODEL),
        tools=[
            agent1_inspect_file,
            agent1_run_python,
            agent1_read_document,
            agent1_calculate_metrics,
            agent1_web_search,
            agent1_crawl_merchant_website,
            agent1_write_log,
        ],
        model_settings=ModelSettings(
            temperature=float(os.getenv("AGENT1_TEMPERATURE", "0.3")),
            max_tokens=int(os.getenv("AGENT1_MAX_TOKENS", "3000")),
        ),
    )


# ---------------------------------------------------------------------------
# Agent 2 — Chat Agent  (minimax/minimax-m2.5)
# Tools: read history directly from output files
# Handoff: to Agent 1 for new research/computation
# ---------------------------------------------------------------------------

AGENT2_MODEL = os.getenv("AGENT2_MODEL", "minimax/minimax-m2.5")

_PDF_FORMAT_INSTRUCTIONS = """
## Quy tắc định dạng khi xuất PDF (export_report_pdf)

Khi user yêu cầu xuất báo cáo PDF, gọi `export_report_pdf` với `analysis_text`.

**CẤU TRÚC GỢI Ý (linh hoạt — bỏ mục không liên quan, thêm mục cần thiết):**
```
## Executive Summary       — tổng TPV, xu hướng chung
## Segment tăng mạnh        — top tăng + driver
## Segment giảm             — top giảm + chẩn đoán nguyên nhân
## Segment mới              — nếu có
## Hành động đề xuất        — cụ thể theo merchant/channel
## Web Research Findings    — bằng chứng + confidence (nếu có)
```

**QUY TẮC FORMAT ĐỂ TRÁNH LỖI PDF:**
- Chỉ dùng: `##` (heading lớn), `###` (heading nhỏ), `- ` (bullet), `| ` (table), `**text**` (bold)
- KHÔNG dùng: emoji trong heading, ký tự đặc biệt Unicode như ₀₁₂→←↑↓✓✗ (gây lỗi font)
- KHÔNG dùng: HTML tags, code blocks (```), subscript/superscript Unicode
- Số liệu: dùng ký tự ASCII (%, +, -, >, <) — KHÔNG dùng ký tự toán học đặc biệt
- Tên merchant tiếng Việt: giữ nguyên dấu (DejaVu font hỗ trợ UTF-8)
- Mỗi section cách nhau bằng dòng trống
- Độ dài tối ưu: 500–2000 từ để PDF không quá dài
"""

AGENT2_INSTRUCTIONS = f"""{ANALYTICS_INSTRUCTIONS}

You are the user-facing Zalopay merchant analytics assistant.

Workflow:
1. **Câu hỏi lịch sử** (TPV tháng trước, lần phân tích trước, trend theo thời gian...):
   - Dùng search_history_files tìm log `agent1_log_*` liên quan đến merchant/period đó
   - Dùng read_history_file đọc nội dung → trả lời trực tiếp, KHÔNG cần gọi Agent 1
   - Có thể so sánh nhiều log để tính trend TPV theo thời gian

2. **Cần tính toán mới** (file mới upload, web research, root-cause chưa có trong log):
   - Handoff sang Research Agent — Agent 1 sẽ tự ghi log sau khi hoàn thành

3. Tổng hợp findings và trả lời user bằng markdown rõ ràng.

4. **Khi user yêu cầu xuất/tải PDF** ("tóm tắt + tải PDF", "xuất báo cáo", "download PDF", "lưu file"...):
   - Soạn nội dung tóm tắt các ý chính (cấu trúc linh hoạt, xem gợi ý bên dưới)
   - GỌI `export_report_pdf(analysis_text=<nội dung>, title=<tiêu đề>)`
   - Sau khi tool trả về, báo "Đã tạo báo cáo PDF, file đang được tải về" (KHÔNG nói path)

5. KHÔNG nhắc đến file path, CSV, JSON, output folder với user.
6. KHÔNG emit tool-call syntax trong câu trả lời cuối.

{_PDF_FORMAT_INSTRUCTIONS}

Data contract that the Research Agent follows:
{DATA_PREPARATION_INSTRUCTIONS}
"""

_NO_TOOL_CALL_INSTRUCTION = (
    "\nCRITICAL: Never output raw tool-call syntax (tool_name:, [TOOL_CALL], function_call, "
    "JSON arguments, or tool names) in your final answer to the user."
)


@function_tool
def search_history_files(query: str = "", limit: int = 6) -> str:
    """Search previously generated analytics reports and agent logs using semantic RAG search.
    Falls back to keyword search if RAG is unavailable.

    Args:
        query: Natural language query (merchant name, channel, period, issue type, etc.).
        limit: Maximum number of results to return (default 6).
    """
    # RAG semantic search
    if _RAG_AVAILABLE and query:
        try:
            chunks = rag.query(query, k=limit)
            if chunks:
                return "\n\n---\n\n".join(chunks)
        except Exception:
            pass
    # Fallback: keyword search trên file system
    return search_output_files(query=query, limit=limit)


@function_tool
def read_history_file(file_path: str, max_characters: int = 12000) -> str:
    """Read the content of a previously generated metrics or report file.

    Args:
        file_path: Path or filename of the output file to read.
        max_characters: Maximum characters to return (default 12000).
    """
    return read_output_file(file_path=file_path, max_characters=max_characters)


@function_tool
async def export_report_pdf(analysis_text: str, title: str = "Báo Cáo Phân Tích Merchant") -> str:
    """Export the current analysis as a PDF report and return the file path for download.
    Call this when the user asks to export or download a report.

    Args:
        analysis_text: The full analysis text to include in the PDF.
        title: Report title (default in Vietnamese with diacritics).
    """
    # Khôi phục dấu tiếng Việt cho cả nội dung lẫn tiêu đề trước khi render PDF
    analysis_text = await restore_vietnamese_diacritics(analysis_text)
    title = await restore_vietnamese_diacritics(title)
    return export_analysis_pdf(analysis_text=analysis_text, title=title)


def _build_agent2(agent1: Agent) -> Agent:
    return Agent(
        name="Analytics Assistant",
        instructions=AGENT2_INSTRUCTIONS + _NO_TOOL_CALL_INSTRUCTION,
        model=_make_client(AGENT2_MODEL),
        tools=[search_history_files, read_history_file, export_report_pdf],
        handoffs=[handoff(agent1)],
        model_settings=ModelSettings(
            temperature=float(os.getenv("AGENT2_TEMPERATURE", "1.0")),
            max_tokens=int(os.getenv("AGENT2_MAX_TOKENS", "2000")),
            top_p=float(os.getenv("AGENT2_TOP_P", "0.95")),
        ),
    )


# ---------------------------------------------------------------------------
# Agent 3 — Background Deep-Analysis Agent  (qwen/qwen3-5-27b)
# Chạy LẶNG THẦM khi user upload file: phân tích đa chiều, phân loại dữ liệu,
# rồi ghi vào RAG để Agent 1/2 query sau này. KHÔNG ảnh hưởng phản hồi chat.
# ---------------------------------------------------------------------------

AGENT3_MODEL = os.getenv("AGENT3_MODEL", "qwen/qwen3-5-27b")

AGENT3_INSTRUCTIONS = """You are Agent 3 — the Zalopay background deep-analysis agent.
You run SILENTLY after a file is uploaded. The user never sees your output directly;
your only job is to enrich the persistent knowledge base (RAG) with deep, multi-dimensional
analysis so Agent 1/2 can answer richer questions later.

## Loại file
- TABULAR (CSV/XLSX/JSON): theo workflow bên dưới (inspect_file → run_python).
- DOCUMENT (PDF/Word/PPT/TXT): dùng `read_document` để lấy nội dung, rồi rút trích & phân loại các
  số liệu/insight quan trọng vào log (không dùng run_python breakdown cho document).

## Workflow (tabular — always follow)
1. `inspect_file` — understand the real structure (sheets, columns, dtypes, ranges). Never assume schema.
2. `run_python` — compute a THOROUGH multi-dimensional breakdown. Be TURN-EFFICIENT: batch
   MANY breakdowns into a FEW run_python calls (aim for 2-4 calls total, not dozens). Cover:
   - Totals & trends per time period (month/quarter) if a date/period exists
   - Breakdown by EVERY categorical dimension found (merchant, channel/SOF, region, acquisition, etc.)
   - Cross-tabs between dimensions (e.g. merchant × channel, channel × period)
   - Growth / change rates period-over-period where possible
   - Ranking: top & bottom performers per dimension
   - Classification of each segment into tiers: High growth / Stable / Underperforming / New / Dropped
   - Anomalies / outliers (sudden spikes or drops); concentration (share of top N in total TPV)
   Print everything you need in those few calls.
3. `write_log` — write ONE comprehensive structured markdown log capturing ALL findings above.
   - task_type = 'deep_analysis'
   - merchants = comma-separated merchant names found in the data
   - Self-contained & queryable: clear headings + actual numbers + classifications.
     Future readers retrieve CHUNKS, so each section must stand alone with concrete numbers.

## Rules (IMPORTANT)
- You have a LIMITED number of turns. Do 2-4 big run_python calls, then ALWAYS call write_log.
- NEVER finish without calling write_log — a partial log is far better than no log.
- Keep the write_log `content` FOCUSED (roughly 400-1200 words): concise headings + key numbers +
  classifications. Do NOT dump raw giant tables — summarize. Over-long content can break the tool call.
- In run_python, print compact summaries (use .head(), rounded numbers) — avoid printing huge frames.
- Never fabricate. Only report what run_python actually returns.
- You do NOT talk to the user and do NOT hand off. After write_log, stop.
"""


def _build_agent3() -> Agent:
    return Agent(
        name="Background Analyst",
        instructions=AGENT3_INSTRUCTIONS,
        model=_make_client(AGENT3_MODEL),
        tools=[agent1_inspect_file, agent1_run_python, agent1_read_document, agent1_write_log],
        model_settings=ModelSettings(
            temperature=float(os.getenv("AGENT3_TEMPERATURE", "0.2")),
            # Lớn để write_log với nội dung dài KHÔNG bị cắt giữa chuỗi JSON (gây 400 Unterminated string)
            max_tokens=int(os.getenv("AGENT3_MAX_TOKENS", "8000")),
        ),
    )


# Giữ reference các task nền để không bị garbage-collected giữa chừng
_BACKGROUND_TASKS: set = set()


async def run_background_analysis_async(file_path: str, user_message: str = "") -> None:
    """Chạy Agent 3 lặng thầm: phân tích sâu file rồi ghi vào RAG. Nuốt mọi lỗi."""
    import logging
    log = logging.getLogger(__name__)
    try:
        agent3 = _build_agent3()
        prompt = (
            f"Uploaded file to deep-analyze: {file_path}\n"
            f"User context (optional): {user_message or 'N/A'}\n"
            "Run the full multi-dimensional analysis workflow and write ONE deep_analysis log."
        )
        max_turns = int(os.getenv("AGENT3_MAX_TURNS", "25"))
        await Runner.run(agent3, input=prompt, max_turns=max_turns)
        log.info("Agent 3 background analysis completed for %s", file_path)
    except Exception as exc:  # KHÔNG bao giờ để lỗi nền ảnh hưởng app
        log.warning("Agent 3 background analysis failed: %s", exc)


def kick_off_background_analysis(file_path: Optional[str], user_message: str = "") -> None:
    """Fire-and-forget Agent 3. Không await → không chặn phản hồi của Agent 1/2."""
    if not file_path:
        return
    try:
        loop = asyncio.get_running_loop()
    except RuntimeError:
        return  # không có event loop đang chạy → bỏ qua (vd gọi từ sync context)
    task = loop.create_task(run_background_analysis_async(file_path, user_message))
    _BACKGROUND_TASKS.add(task)
    task.add_done_callback(_BACKGROUND_TASKS.discard)


# ---------------------------------------------------------------------------
# Orchestrator
# ---------------------------------------------------------------------------

def _rag_startup_sync() -> None:
    """Index toàn bộ file output cũ vào ChromaDB lúc startup (chạy 1 lần)."""
    if not _RAG_AVAILABLE:
        return
    try:
        n = rag.index_directory(OUTPUT_DIR, pattern="*.md")
        if n:
            import logging
            logging.getLogger(__name__).info(f"RAG startup: indexed {n} existing output files")
    except Exception:
        pass


_rag_startup_sync()  # chạy khi module được import lần đầu


# Phát hiện output bị leak tool-call / JSON fragment (vd model xuất '}' hoặc '{...}')
_LEAK_PATTERN = re.compile(
    r'^[\s`]*[\{\}\[\]]'                       # bắt đầu bằng { } [ ]
    r'|"(name|arguments|tool_call|function|parameters)"\s*:'  # khóa JSON tool-call
    r'|\[TOOL_CALL\]|<tool_call>|functions\.',
    re.IGNORECASE,
)


def _looks_like_tool_leak(text: str) -> bool:
    s = (text or "").strip()
    if not s:
        return True
    # Quá ngắn và toàn ký tự JSON → chắc chắn leak
    if len(s) <= 3 and s.strip("{}[]() \t\n`"):
        return True
    if s in ("{}", "{", "}", "[]", "[", "]"):
        return True
    return bool(_LEAK_PATTERN.match(s))


def _build_history_input(history: Optional[list], current_message: str) -> Any:
    """Chuyển history [(user, assistant), ...] thành list input items cho SDK + message hiện tại.

    Trả về string nếu không có history (giữ fast-path), ngược lại list message items.
    """
    if not history:
        return current_message
    items: List[Dict[str, str]] = []
    for turn in history[-6:]:  # giữ tối đa 6 lượt gần nhất để không phình context
        try:
            user_msg, bot_msg = turn[0], turn[1]
        except (IndexError, TypeError):
            continue
        if user_msg:
            items.append({"role": "user", "content": str(user_msg)})
        if bot_msg and not _looks_like_tool_leak(str(bot_msg)):
            items.append({"role": "assistant", "content": str(bot_msg)})
    items.append({"role": "user", "content": current_message})
    return items


_LEAK_FALLBACK = (
    "Xin lỗi, tôi chưa tạo được câu trả lời hoàn chỉnh cho yêu cầu này. "
    "Bạn thử diễn đạt lại hoặc nêu rõ tên merchant cần kiểm tra giúp tôi nhé."
)


def _is_transient_model_error(exc: Exception) -> bool:
    """Lỗi tạm thời do model sinh tool-call JSON hỏng / quá tải → nên retry."""
    msg = str(exc).lower()
    return any(
        s in msg for s in
        ("unterminated string", "bad request", "400", "json", "rate limit",
         "429", "timeout", "temporarily", "503", "502", "overloaded")
    )


def _friendly_error_message(exc: Exception) -> str:
    """Thông báo lỗi thân thiện (không dán raw JSON/stacktrace cho user)."""
    if _is_transient_model_error(exc):
        return (
            "Xin lỗi, mô hình gặp trục trặc tạm thời khi xử lý yêu cầu này (thường do sinh dữ liệu "
            "chưa hoàn chỉnh). Bạn vui lòng thử lại — hoặc diễn đạt ngắn gọn/rõ ràng hơn một chút."
        )
    return (
        "Đã xảy ra lỗi khi xử lý yêu cầu. Vui lòng thử lại hoặc kiểm tra kết nối GreenNode."
    )


def _prepare_agent2_input(
    user_message: str,
    uploaded_file_path: Optional[str],
    history: Optional[list],
) -> Tuple["Agent", Any]:
    """Setup chung cho cả chạy thường lẫn streaming: build agents + compose input."""
    _greennode_config()  # validate credentials early
    agent1 = _build_agent1()
    agent2 = _build_agent2(agent1)

    if uploaded_file_path:
        # Kích hoạt Agent 3 chạy NỀN (không await) — phân tích sâu + đẩy RAG, không chặn chat
        kick_off_background_analysis(uploaded_file_path, user_message)
        current = (
            f"{user_message}\n\n"
            f"[Uploaded file: {uploaded_file_path}] "
            "Please hand off to the Research Agent to process this file and then analyze the results."
        )
    else:
        current = user_message

    if not uploaded_file_path and _is_simple_message(user_message) and not history:
        agent_input: Any = current
    else:
        agent_input = _build_history_input(history, current)
    return agent2, agent_input


async def run_merchant_workflow_async(
    user_message: str,
    uploaded_file_path: Optional[str] = None,
    history: Optional[list] = None,
) -> str:
    """Run the multi-agent workflow (non-streaming): Agent 2 orchestrates, Agent 1 researches.

    Args:
        user_message: tin nhắn hiện tại của user.
        uploaded_file_path: file upload (nếu có) → kích hoạt Agent 1 + Agent 3.
        history: list [(user, assistant), ...] các lượt TRƯỚC để giữ ngữ cảnh follow-up.
    """
    user_message = (user_message or "").strip() or "Xin chào!"
    agent2, agent_input = _prepare_agent2_input(user_message, uploaded_file_path, history)
    run_max_turns = int(os.getenv("AGENT_MAX_TURNS", "16"))
    last_exc: Optional[Exception] = None
    for attempt in range(int(os.getenv("FALLBACK_RETRIES", "3"))):  # malformed-JSON ngẫu nhiên → retry
        try:
            result = await Runner.run(agent2, input=agent_input, max_turns=run_max_turns)
            output = (result.final_output or "").strip()
            if _looks_like_tool_leak(output):
                return _LEAK_FALLBACK
            return output or "Không có kết quả. Vui lòng thử lại."
        except Exception as exc:
            last_exc = exc
            if not _is_transient_model_error(exc):
                break
    return _friendly_error_message(last_exc) if last_exc else "Không có kết quả. Vui lòng thử lại."


async def run_merchant_workflow_stream(
    user_message: str,
    uploaded_file_path: Optional[str] = None,
    history: Optional[list] = None,
):
    """Streaming version: yield (partial_text, is_final) khi agent đang sinh chữ.

    - Trong lúc tool-call/handoff (chưa có chữ) → không yield → UI vẫn hiện typing dots.
    - Khi agent sinh text → yield text tích lũy dần (is_final=False).
    - Kết thúc → yield (final_output đã làm sạch, True).
    """
    user_message = (user_message or "").strip() or "Xin chào!"
    try:
        agent2, agent_input = _prepare_agent2_input(user_message, uploaded_file_path, history)
    except Exception as exc:
        yield (f"Đã xảy ra lỗi khi xử lý yêu cầu: {exc}", True)
        return

    # Chỉ stream TEXT đầu ra thật, KHÔNG stream reasoning (ResponseReasoningTextDeltaEvent)
    try:
        from openai.types.responses import ResponseTextDeltaEvent
    except Exception:
        ResponseTextDeltaEvent = None  # type: ignore

    last_exc: Optional[Exception] = None
    accumulated = ""
    run_max_turns = int(os.getenv("AGENT_MAX_TURNS", "16"))  # flow nhiều tool (RAG+web+crawl) cần nhiều lượt
    if True:  # ── Lượt 1: STREAMING (UX nhanh) ──
        try:
            result = Runner.run_streamed(agent2, input=agent_input, max_turns=run_max_turns)
            async for event in result.stream_events():
                if getattr(event, "type", "") != "raw_response_event":
                    continue
                data = getattr(event, "data", None)
                if ResponseTextDeltaEvent is not None:
                    if not isinstance(data, ResponseTextDeltaEvent):
                        continue
                elif type(data).__name__ != "ResponseTextDeltaEvent":
                    continue
                delta = getattr(data, "delta", None)
                if isinstance(delta, str) and delta:
                    accumulated += delta
                    yield (accumulated, False)

            final = (getattr(result, "final_output", None) or accumulated or "").strip()
            if _looks_like_tool_leak(final):
                final = _LEAK_FALLBACK
            yield (final or "Không có kết quả. Vui lòng thử lại.", True)
            return
        except Exception as exc:
            last_exc = exc

    # ── Lượt 2+: FALLBACK NON-STREAMING, retry NHIỀU lần ──
    # Lỗi malformed-JSON là NGẪU NHIÊN → mỗi lần thử lại có cơ hội thành công. Non-streaming để
    # server tự ráp response (ổn hơn streaming). Vẫn fallback dù đã stream preamble ("Để tôi tìm...")
    # vì preamble KHÔNG phải câu trả lời thật — kết quả non-streaming sẽ thay thế nó.
    if _is_transient_model_error(last_exc):
        for _ in range(int(os.getenv("FALLBACK_RETRIES", "3"))):
            try:
                result = await Runner.run(agent2, input=agent_input, max_turns=run_max_turns)
                final = (getattr(result, "final_output", None) or "").strip()
                if _looks_like_tool_leak(final):
                    final = _LEAK_FALLBACK
                yield (final or "Không có kết quả. Vui lòng thử lại.", True)
                return
            except Exception as exc2:
                last_exc = exc2
                if not _is_transient_model_error(exc2):
                    break  # lỗi không phải tạm thời → ngừng retry

    yield (_friendly_error_message(last_exc) if last_exc else "Không có kết quả. Vui lòng thử lại.", True)


def run_merchant_workflow(user_message: str, uploaded_file_path: Optional[str] = None) -> str:
    """Synchronous wrapper for scripts and simple callers."""
    return asyncio.run(run_merchant_workflow_async(user_message, uploaded_file_path))
